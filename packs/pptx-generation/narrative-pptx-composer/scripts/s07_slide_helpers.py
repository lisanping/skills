"""
Shared slide-building helpers for narrative-pptx-composer.

Imported by per-session `s07-build.py` scripts to compose slides
with python-pptx. Provides reusable building blocks for text
boxes, cards, accent bars, images, rounded rects, etc.
"""
import os
from pathlib import Path

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


# ── Color helpers ──────────────────────────────────────────────────

def resolve_color(value):
    """Convert a hex string (e.g. '0078D4') to RGBColor.

    Accepts:
      - 6-char hex string: '0078D4'
      - RGBColor instance (returned as-is)
      - Tuple of (r, g, b) ints
    """
    if isinstance(value, RGBColor):
        return value
    if isinstance(value, (list, tuple)) and len(value) == 3:
        return RGBColor(*value)
    if isinstance(value, str):
        v = value.lstrip("#")
        return RGBColor(int(v[0:2], 16), int(v[2:4], 16), int(v[4:6], 16))
    raise ValueError(f"Cannot resolve color: {value!r}")


# ── Font / size helpers ────────────────────────────────────────────

def resolve_font(font_token, typography):
    """Resolve a font token ('heading', 'body', 'mono') to a font name."""
    if font_token == "heading":
        return typography.get("headingFont", "Segoe UI Semibold")
    if font_token == "mono":
        return typography.get("monoFont", "JetBrains Mono")
    return typography.get("bodyFont", "Segoe UI")


def resolve_size(size_token, typography):
    """Resolve a size token or raw number to Pt value."""
    if isinstance(size_token, (int, float)):
        return Pt(size_token)
    scale = typography.get("scale", {})
    pt_val = scale.get(size_token, scale.get("body", 14))
    return Pt(pt_val)


# ── Slide background ──────────────────────────────────────────────

def set_slide_bg(slide, color):
    """Set a solid background color on a slide.

    Args:
        slide: pptx slide object.
        color: RGBColor instance.
    """
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _apply_shape_alpha(shape, alpha):
    """Apply an alpha (0.0–1.0) to a solid-filled shape via direct XML."""
    if not (alpha and alpha > 0 and alpha < 1):
        return
    from pptx.oxml.ns import qn
    from lxml import etree
    sp_pr = shape.fill._xPr
    solid_fill = sp_pr.find(qn("a:solidFill"))
    if solid_fill is None:
        return
    srgb = solid_fill.find(qn("a:srgbClr"))
    if srgb is None:
        return
    alpha_val = str(int(max(0.0, min(1.0, alpha)) * 100000))
    alpha_el = etree.SubElement(srgb, qn("a:alpha"))
    alpha_el.set("val", alpha_val)


def set_slide_bg_image(slide, img_path, *, fit="cover",
                       fallback_color=None, slide_size=None,
                       scrim_color=None, scrim_alpha=0.0):
    """Lay a full-bleed background image behind all other shapes.

    Inserts the image at slide (0, 0) sized to **cover** the full
    slide (fills the canvas, may crop overflow), then sends it to
    the back so subsequent shapes paint on top. Optionally paints a
    full-bleed scrim rectangle on top of the image for legibility of
    overlaid text — the scrim is the LAST element added before this
    function returns, so any text added by the caller afterwards
    paints above it.

    If `img_path` is None / empty / missing, falls back to
    `set_slide_bg(slide, fallback_color)` (a solid color) so callers
    can degrade gracefully when image generation failed. The scrim
    is NOT painted in the fallback path — the caller should pick a
    `fallback_color` that already provides text contrast.

    Args:
        slide: pptx slide object.
        img_path: path to the background image (str or Path).
        fit: "cover" (default; fill, may crop) — recommended for true
             full-bleed backgrounds. "contain" (fit, may letterbox)
             only when preserving the entire image is required; the
             letterbox bands will show the slide's existing background
             color, so prefer setting that with `set_slide_bg` first.
        fallback_color: RGBColor used when the image is unavailable.
                        Defaults to white.
        slide_size: optional (width_emu, height_emu) tuple. If omitted,
                    derived from `slide.part.package.presentation_part
                    .presentation.slide_width / slide_height`.
        scrim_color: optional RGBColor for a full-bleed scrim overlay.
                     None disables the scrim. Pass black for darkening
                     under light text; white for lightening under dark
                     text. Decision should be driven per-slide by
                     headline color, image luminance, and visual tone
                     — do not hardcode for the deck.
        scrim_alpha: 0.0–1.0 opacity. 0.0 (default) means no scrim.
                     Typical legibility values: 0.25–0.45 for editorial
                     overlays; 0.55–0.75 for hero covers with bold text
                     over busy imagery.

    Returns:
        (picture_shape, scrim_shape) — either may be None.
        picture_shape is None when the fallback color path was used;
        scrim_shape is None when no scrim was requested or the
        fallback path was taken.
    """
    if fallback_color is None:
        fallback_color = RGBColor(0xFF, 0xFF, 0xFF)

    if not img_path:
        set_slide_bg(slide, fallback_color)
        return None, None

    img_path = Path(img_path) if not isinstance(img_path, Path) else img_path
    if not img_path.exists():
        set_slide_bg(slide, fallback_color)
        return None, None

    if slide_size is None:
        prs = slide.part.package.presentation_part.presentation
        slide_w, slide_h = prs.slide_width, prs.slide_height
    else:
        slide_w, slide_h = slide_size

    try:
        from PIL import Image
        img = Image.open(img_path)
        img_w, img_h = img.size
        img_aspect = img_w / img_h
        slide_aspect = slide_w / slide_h

        if fit == "contain":
            if img_aspect > slide_aspect:
                disp_w = slide_w
                disp_h = int(slide_w / img_aspect)
            else:
                disp_h = slide_h
                disp_w = int(slide_h * img_aspect)
            left = (slide_w - disp_w) // 2
            top = (slide_h - disp_h) // 2
        else:  # cover
            if img_aspect > slide_aspect:
                disp_h = slide_h
                disp_w = int(slide_h * img_aspect)
            else:
                disp_w = slide_w
                disp_h = int(slide_w / img_aspect)
            left = (slide_w - disp_w) // 2
            top = (slide_h - disp_h) // 2
    except ImportError:
        left, top = 0, 0
        disp_w, disp_h = slide_w, slide_h

    pic = slide.shapes.add_picture(
        str(img_path), left, top, disp_w, disp_h)

    # Send picture to the back so other shapes render on top.
    spTree = pic._element.getparent()
    spTree.remove(pic._element)
    spTree.insert(2, pic._element)

    scrim = None
    if scrim_color is not None and scrim_alpha and scrim_alpha > 0:
        scrim = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, slide_w, slide_h)
        scrim.fill.solid()
        scrim.fill.fore_color.rgb = scrim_color
        scrim.line.fill.background()
        _apply_shape_alpha(scrim, scrim_alpha)

    return pic, scrim


# ── Text box ──────────────────────────────────────────────────────

def add_textbox(slide, left, top, width, height, text,
                font_size=18, bold=False, color=None,
                alignment=PP_ALIGN.LEFT, font_name="Segoe UI",
                anchor=MSO_ANCHOR.TOP, line_spacing=1.15):
    """Add a text box with multi-paragraph support.

    Splits text on newlines and applies uniform formatting.
    Returns the textbox shape.
    """
    if color is None:
        color = RGBColor(0x23, 0x23, 0x23)

    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    try:
        tf.vertical_anchor = anchor
    except Exception:
        pass

    paragraphs = text.split("\n")
    for i, para_text in enumerate(paragraphs):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = para_text
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.font.name = font_name
        p.alignment = alignment
        p.space_after = Pt(4)
        p.line_spacing = Pt(font_size * line_spacing)

    return txBox


# ── Text box (autofit-to-zone) ────────────────────────────────────

def _estimate_lines(text, width_emu, font_size_pt, bold=False):
    """Rough line-count estimate for `text` wrapping inside `width_emu`
    at `font_size_pt`. Uses a conservative average glyph width
    (0.55 × font_size for Latin proportional fonts; 0.62 if bold).
    Counts explicit newlines, then wraps each line by character budget.
    """
    if not text:
        return 1
    # 914400 EMU per inch; PowerPoint default = 96 DPI; 1 pt = 1/72 in.
    # px-equivalent of font advance ≈ 0.55 * pt * (96/72) = 0.733 * pt
    # But we want "characters per inch" to compare to width in inches.
    # avg_char_width_in_pt ≈ 0.55 * font_size  (a Latin heuristic)
    # 1 pt = 1/72 inch → avg_char_width_in_inch = (0.55 * fs) / 72
    factor = 0.62 if bold else 0.55
    width_in = width_emu / 914400.0
    char_width_in = (factor * font_size_pt) / 72.0
    if char_width_in <= 0:
        return 1
    chars_per_line = max(1, int(width_in / char_width_in))
    total_visual_lines = 0
    for raw in text.split("\n"):
        if not raw.strip():
            total_visual_lines += 1
            continue
        # Word-aware wrap estimate: ceil(len / chars_per_line) is too
        # generous; bump by 1 on long unwrappable runs.
        n = len(raw)
        total_visual_lines += max(1, -(-n // chars_per_line))  # ceil division
    return total_visual_lines


def add_textbox_autofit(slide, left, top, width, height, text,
                        font_size=24, bold=False, color=None,
                        alignment=PP_ALIGN.LEFT, font_name="Segoe UI",
                        anchor=MSO_ANCHOR.TOP, line_spacing=1.15,
                        min_size=12, step=2):
    """Escape-hatch only. Step 6b owns content fit; Step 7 builders
    should always use `add_textbox`.

    Like `add_textbox`, but steps `font_size` down (by `step` pt)
    until the estimated wrapped line count fits inside `height`.
    Stops at `min_size` even if the text still overflows.

    Calling this from `s07-build.py` is a signal that Step 6b's
    wrap-estimate heuristic mis-judged this slot — log the trigger
    in `s09-session-retrospective.json` so the heuristic can be
    tuned, do not normalise it as a routine Step 7 tool.

    Returns the created textbox shape.
    """
    fs = float(font_size)
    # Each wrapped visual line takes ≈ font_size * line_spacing points
    # of vertical space. Convert height (EMU) → points: 1 pt = 12700 EMU.
    height_pt = height / 12700.0
    while fs >= min_size:
        line_count = _estimate_lines(text, width, fs, bold=bold)
        used_pt = line_count * fs * line_spacing
        if used_pt <= height_pt:
            break
        fs -= step
    fs = max(min_size, fs)
    return add_textbox(slide, left, top, width, height, text,
                       font_size=fs, bold=bold, color=color,
                       alignment=alignment, font_name=font_name,
                       anchor=anchor, line_spacing=line_spacing)


# ── Accent bar ────────────────────────────────────────────────────

def add_accent_bar(slide, left, top,
                   width=Inches(1.2), height=Inches(0.06),
                   color=None):
    """Add a thin decorative accent bar."""
    if color is None:
        color = RGBColor(0x00, 0x78, 0xD4)
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


# ── Outline-only shapes (no fill) ─────────────────────────────────

def add_outline_capsule(slide, left, top, width, height,
                        outline_color, outline_weight_pt=2.0):
    """Outline-only rounded rectangle. With width >> height and full
    corner adjustment, renders as a capsule/pill. Use for editorial
    "outline-only chip" treatments where colored borders carry the
    visual weight (no fill blocks compete with text)."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.background()
    shape.line.color.rgb = resolve_color(outline_color)
    shape.line.width = Pt(outline_weight_pt)
    if hasattr(shape, "adjustments") and len(shape.adjustments) > 0:
        shape.adjustments[0] = 0.5  # full pill / max corner radius
    return shape


def add_outline_circle(slide, left, top, width, height,
                       outline_color, outline_weight_pt=2.0):
    """Outline-only ellipse/circle. Pair with an icon textbox or
    image overlay to build "icon-in-colored-ring" badges."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, left, top, width, height
    )
    shape.fill.background()
    shape.line.color.rgb = resolve_color(outline_color)
    shape.line.width = Pt(outline_weight_pt)
    return shape


# ── Image (safe) ──────────────────────────────────────────────────

def add_image_safe(slide, img_path, left, top, width, height=None,
                   placeholder_color=None):
    """Add an image preserving aspect ratio, with missing-file fallback.

    If `height` is provided, computes display dimensions that fit
    within (width × height) while preserving the original aspect
    ratio. If `img_path` is None / empty / missing, renders a grey
    placeholder rectangle (lets the caller's `fallback` reflow run
    elsewhere without raising). Returns the created shape, or None
    if `img_path` is None / empty.
    """
    if not img_path:
        return None

    if placeholder_color is None:
        placeholder_color = RGBColor(0xE8, 0xEB, 0xF0)

    img_path = Path(img_path) if not isinstance(img_path, Path) else img_path
    if not img_path.exists():
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, left, top, width,
            height or Inches(3)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = placeholder_color
        add_textbox(slide, left + Inches(0.3), top + Inches(0.3),
                    width - Inches(0.6), Inches(0.5),
                    f"[Image: {img_path.name}]", font_size=11,
                    color=RGBColor(0x55, 0x55, 0x55))
        return shape

    if height:
        try:
            from PIL import Image
            img = Image.open(img_path)
            img_w, img_h = img.size
            aspect = img_w / img_h

            # Convert Emu to float for ratio calculation
            zone_w = width if isinstance(width, (int, float)) else width
            zone_h = height if isinstance(height, (int, float)) else height
            zone_aspect = zone_w / zone_h

            if aspect > zone_aspect:
                display_w = zone_w
                display_h = int(zone_w / aspect)
            else:
                display_h = zone_h
                display_w = int(zone_h * aspect)
            return slide.shapes.add_picture(
                str(img_path), left, top, display_w, display_h)
        except ImportError:
            # Pillow not available — fall back to explicit dimensions
            return slide.shapes.add_picture(
                str(img_path), left, top, width, height)
    else:
        return slide.shapes.add_picture(
            str(img_path), left, top, width=width)


# ── Image with overlay (image + scrim for legible text on top) ────

def add_image_with_overlay(slide, img_path, left, top, width, height,
                           *, scrim_color=None, scrim_alpha=0.4,
                           fallback_color=None):
    """Place an image inside a zone with a semi-transparent scrim on top.

    Use when text needs to sit legibly over a photographic / generated
    image. Caller is responsible for adding the text afterwards (the
    text will paint above the scrim because it is added later).

    If `img_path` is None / missing, paints `fallback_color` as a solid
    rectangle so the zone retains its compositional weight.

    Args:
        slide: pptx slide object.
        img_path: path to image (str / Path) or None for fallback.
        left, top, width, height: zone geometry in EMU (use Inches()).
        scrim_color: RGBColor for the overlay. None disables the scrim.
                     Default: black for darkening behind light text.
        scrim_alpha: 0.0–1.0 opacity. 0.4 is a sane default for white
                     text over a busy image.
        fallback_color: solid color used when image is unavailable.

    Returns:
        (image_shape_or_None, scrim_shape_or_None)
    """
    if scrim_color is None:
        scrim_color = RGBColor(0x00, 0x00, 0x00)
    if fallback_color is None:
        fallback_color = RGBColor(0x33, 0x33, 0x33)

    img_shape = None
    if img_path:
        img_shape = add_image_safe(slide, img_path, left, top, width, height)

    if img_shape is None and not img_path:
        # No image at all — paint a solid block so the zone still reads.
        img_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, left, top, width, height)
        img_shape.fill.solid()
        img_shape.fill.fore_color.rgb = fallback_color
        img_shape.line.fill.background()

    scrim = None
    if scrim_alpha and scrim_alpha > 0:
        scrim = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, left, top, width, height)
        scrim.fill.solid()
        scrim.fill.fore_color.rgb = scrim_color
        scrim.line.fill.background()
        _apply_shape_alpha(scrim, scrim_alpha)

    return img_shape, scrim


# ── Rounded rectangle ────────────────────────────────────────────

def add_rounded_rect(slide, left, top, width, height, fill_color,
                     text="", font_size=14, text_color=None,
                     bold=False, corner_radius_in=0.12,
                     corner_radius=None):
    """Add a rounded rectangle, optionally with centered text.

    `corner_radius_in` is an **absolute radius in inches** (default
    0.12in ≈ 8.6pt — modern card radius). The helper internally
    converts to python-pptx's fraction-of-shorter-side adjustment.
    Cap is 0.5 of shorter side (full pill).

    `corner_radius` (legacy positional arg, fraction 0.0–0.5) is
    still accepted for back-compat but a deprecation note is logged.
    Prefer `corner_radius_in` for new code so radii stay visually
    consistent across small and large cards (a 0.15 fraction on a
    7-inch-tall card = 1.05in radius — visually overrounded).
    """
    if text_color is None:
        text_color = RGBColor(0xFF, 0xFF, 0xFF)

    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    if hasattr(shape, 'adjustments') and len(shape.adjustments) > 0:
        if corner_radius is not None:
            # Legacy fraction path
            shape.adjustments[0] = max(0.0, min(0.5, corner_radius))
        else:
            # Convert absolute inches → fraction of shorter side.
            shorter_emu = min(width, height)
            shorter_in = shorter_emu / 914400.0
            if shorter_in > 0:
                fraction = corner_radius_in / shorter_in
                shape.adjustments[0] = max(0.0, min(0.5, fraction))

    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = text_color
        p.font.bold = bold
        p.font.name = "Segoe UI"
        p.alignment = PP_ALIGN.CENTER

    return shape


# ── Card (compound element) ──────────────────────────────────────

def add_card(slide, left, top, width, height,
             title_text, body_text, num_text="",
             accent_color=None, card_bg=None,
             border_color=None, title_size=14, body_size=11,
             corner_radius_in=0.12):
    """Add a card with optional number badge, title, and body text.

    Composes: rounded rect background + number circle + title
    textbox + body textbox.

    `corner_radius_in` is the absolute corner radius in inches
    (default 0.12in ≈ 8.6pt). Internally converted to python-pptx's
    fraction-of-shorter-side adjustment so the visual radius stays
    constant regardless of card size — without this, a tall card
    (e.g. 3×6in) at the previous hard-coded 0.06 fraction got a
    0.36in radius vs. a small card (3×2in) which got 0.12in,
    making the deck look inconsistent.
    """
    if accent_color is None:
        accent_color = RGBColor(0x00, 0x78, 0xD4)
    if card_bg is None:
        card_bg = RGBColor(0xFF, 0xFF, 0xFF)
    if border_color is None:
        border_color = RGBColor(0xDD, 0xDD, 0xDD)

    # Card background
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    card.fill.solid()
    card.fill.fore_color.rgb = card_bg
    card.line.color.rgb = border_color
    card.line.width = Pt(1)
    if hasattr(card, 'adjustments') and len(card.adjustments) > 0:
        shorter_in = min(width, height) / 914400.0
        if shorter_in > 0:
            card.adjustments[0] = max(0.0, min(0.5, corner_radius_in / shorter_in))

    # Number badge
    if num_text:
        badge = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, left + Inches(0.2), top + Inches(0.2),
            Inches(0.45), Inches(0.45)
        )
        badge.fill.solid()
        badge.fill.fore_color.rgb = accent_color
        badge.line.fill.background()
        tf = badge.text_frame
        p = tf.paragraphs[0]
        p.text = num_text
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        p.font.bold = True
        p.font.name = "Segoe UI"
        p.alignment = PP_ALIGN.CENTER

    # Title
    add_textbox(slide, left + Inches(0.2), top + Inches(0.75),
                width - Inches(0.4), Inches(0.4),
                title_text, font_size=title_size, bold=True,
                color=RGBColor(0x23, 0x23, 0x23))

    # Body
    add_textbox(slide, left + Inches(0.2), top + Inches(1.15),
                width - Inches(0.4), height - Inches(1.35),
                body_text, font_size=body_size,
                color=RGBColor(0x55, 0x55, 0x55),
                line_spacing=1.25)

    return card


# ── Speaker notes ─────────────────────────────────────────────────

def add_speaker_notes(slide, text):
    """Add speaker notes to a slide."""
    if text:
        notes_slide = slide.notes_slide
        tf = notes_slide.notes_text_frame
        tf.text = text


# ── Slot lookup ───────────────────────────────────────────────────

def get_slot(slots, zone):
    """Find a slot by zone name in a slots list. Returns None if missing."""
    for s in slots:
        if s.get("zone") == zone:
            return s
    return None


def resolve_image_path(plan_zone, slot=None):
    """Resolve an image-zone path with the mandatory plan-first precedence.

    After Step 5f patches, the resolved image path lives on the **plan
    zone** (`s05-slide-visual-design.json → slides[i].layoutSpec.zones[j]
    .path`), NOT on the content slot. `s06-slide-content.json` slot
    entries for image zones may carry `"path": null` by design.

    Reading only from the slot silently produces empty placeholders for
    every generated illustration — invisible to text-based validators
    (`s06_validate_content.py` only checks slot presence) and only
    caught by Step 8 visual QA after the deck is rendered. Always
    resolve image paths through this helper.

    Args:
        plan_zone: zone dict from `s05-slide-visual-design.json`
                   (`slides[i].layoutSpec.zones[j]`). Required.
        slot:      optional slot dict from `s06-slide-content.json`
                   (`slides[i].slots[k]`). Used as a fallback only when
                   the plan zone has no `path`.

    Returns:
        str path if found on either side, else None. Pass the result
        directly to `add_image_safe` / `add_image_with_overlay` /
        `set_slide_bg_image` — those helpers already handle a None
        path via their fallback paths.
    """
    plan_path = (plan_zone or {}).get("path")
    if plan_path:
        return plan_path
    if slot:
        return slot.get("path")
    return None


# ── Title-choreography helpers (Phase A — see workflow-step5 Rule 4) ──

def add_textbox_runs(slide, left, top, width, height, runs,
                     base_font_name="Segoe UI", base_color=None,
                     base_size=24, alignment=PP_ALIGN.LEFT,
                     anchor=MSO_ANCHOR.TOP, line_spacing=1.15,
                     paragraph_spacing_pt=4):
    """Add a textbox composed of mixed-style runs in a single paragraph.

    Implements the **mixed-weight-run** title move: a single title
    can weave together regular and bold (and color/italic) runs to
    encode structural emphasis inside the headline itself, e.g.
    "AI changes **tasks** before it changes jobs".

    `runs` is a list of dicts. Each run renders as a separate
    `<a:r>` inside the same paragraph (no line break between runs).
    Use an explicit run with `text: "\\n"` to start a new paragraph.

    Run schema (all keys optional except `text`):
        {
            "text":   "AI changes ",          # required
            "bold":   True,                   # default False
            "italic": False,                  # default False
            "color":  RGBColor(...),          # default base_color
            "font":   "DM Sans",              # default base_font_name
            "size":   38,                     # pt; default base_size
        }

    Args:
        slide:           pptx slide object.
        left/top/width/height: EMU coordinates (use Inches()).
        runs:            list of run dicts (see schema above).
        base_font_name:  fallback font when a run omits `font`.
        base_color:      fallback RGBColor when a run omits `color`.
        base_size:       fallback point size when a run omits `size`.
        alignment:       paragraph alignment.
        anchor:          vertical anchor (top/middle/bottom).
        line_spacing:    multiplier of base_size for line height.
        paragraph_spacing_pt: extra space-after on the paragraph.

    Returns:
        the textbox shape.

    Notes:
        * No autofit. Caller picks a `base_size` that fits — this
          helper is meant for editorial title moments, not bulk text.
        * To start a new paragraph, insert a run dict where the only
          key is `{"text": "\\n"}` and the next run will land on a
          fresh line with the same default formatting.
    """
    if base_color is None:
        base_color = RGBColor(0x23, 0x23, 0x23)

    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    try:
        tf.vertical_anchor = anchor
    except Exception:
        pass

    # First paragraph (always exists in a fresh text_frame).
    p = tf.paragraphs[0]
    p.alignment = alignment
    p.space_after = Pt(paragraph_spacing_pt)
    p.line_spacing = Pt(base_size * line_spacing)

    # Strip the empty default run so we control formatting per added run.
    # python-pptx exposes p.runs via the underlying lxml element; we
    # add runs explicitly rather than mutating the default empty one.
    for r in list(p.runs):
        r._r.getparent().remove(r._r)

    def _emit(p_target, run_spec):
        text = run_spec.get("text", "")
        if text == "\n":
            # Caller asked for a new paragraph — return a fresh one.
            new_p = tf.add_paragraph()
            new_p.alignment = alignment
            new_p.space_after = Pt(paragraph_spacing_pt)
            new_p.line_spacing = Pt(base_size * line_spacing)
            return new_p
        run = p_target.add_run()
        run.text = text
        run.font.name = run_spec.get("font", base_font_name)
        size_pt = run_spec.get("size", base_size)
        run.font.size = Pt(size_pt)
        run.font.bold = bool(run_spec.get("bold", False))
        run.font.italic = bool(run_spec.get("italic", False))
        run.font.color.rgb = run_spec.get("color", base_color)
        return p_target

    current_p = p
    for run_spec in runs:
        result = _emit(current_p, run_spec)
        if result is not current_p:
            current_p = result

    return txBox


def add_oversized_numeral(slide, left, top, width, height, char,
                          color=None, font_name="Segoe UI",
                          font_size_pt=160, bold=True,
                          alignment=PP_ALIGN.LEFT,
                          anchor=MSO_ANCHOR.TOP):
    """Add a single oversized character (numeral, letter, or symbol)
    sized to anchor a section opener or stat callout.

    Implements the **oversized-numeral** title move. The character
    is rendered in a textbox at `font_size_pt` with no autofit.
    Recommended sizing: width 1.5–3.0in, height 1.0–2.0in,
    font_size 100–180pt for a 1–3 character glyph. Zone area
    around 4–10% of a 16:9 canvas reads as editorial scale
    (the validator's `oversized-numeral` recognition signal
    requires >= 4% canvas area).

    Args:
        slide:        pptx slide object.
        left/top/width/height: EMU coordinates (use Inches()).
        char:         single character (or short token like "01", "Q3").
        color:        RGBColor; defaults to dark grey.
        font_name:    typeface; pair with deck heading font.
        font_size_pt: point size; values < 80pt usually do not register
                      as oversized — prefer 120–200pt for true editorial scale.
        bold:         default True.
        alignment:    horizontal alignment within the box.
        anchor:       vertical anchor.

    Returns:
        the textbox shape.

    Notes:
        * Place this zone with `role: "numeral"` in layoutSpec.zones
          so the title-choreography validator can recognize it as
          satisfying the `oversized-numeral` move.
        * Set `text` field in s06-slide-content.json slot to the
          character. A 1–3 char value is recommended.
    """
    if color is None:
        color = RGBColor(0x23, 0x23, 0x23)

    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = False
    tf.auto_size = None
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    try:
        tf.vertical_anchor = anchor
    except Exception:
        pass

    p = tf.paragraphs[0]
    p.text = char
    p.alignment = alignment
    p.font.name = font_name
    p.font.size = Pt(font_size_pt)
    p.font.bold = bool(bold)
    p.font.color.rgb = color
    # Tight line spacing so the glyph hugs the box.
    p.line_spacing = 1.0

    return txBox


def add_vertical_label(slide, left, top, width, height, text,
                       color=None, font_name="Segoe UI",
                       font_size_pt=10, bold=True, letter_spacing_em=0.18,
                       direction="vert270"):
    """Add a small vertically-set label (typically along a slide edge).

    Implements the **vertical-act-label** title move: a thin running
    label such as "ACT 02 · MECHANISM" stacked along the left or
    right edge of the slide, providing a deck-wide structural
    signal without intruding on the title area.

    Implementation uses PowerPoint's `bodyPr` `vert` attribute on
    the text frame — this rotates the text rendering, not the shape
    bounding box, so the textbox `width`/`height` you pass in match
    the rotated visual footprint. For a left-side label running
    bottom-to-top, use direction "vert270" with a tall, narrow box
    (e.g. width=0.4in, height=4.5in).

    Args:
        slide:            pptx slide object.
        left/top/width/height: EMU coordinates of the rotated box
                          (NOT the unrotated text). Use Inches().
        text:             the label string. Will be uppercased and
                          letter-spaced for editorial feel.
        color:            RGBColor; defaults to muted grey.
        font_name:        typeface; pair with deck heading font.
        font_size_pt:     point size; 8–12pt is typical.
        bold:             default True for stronger "label" feel.
        letter_spacing_em: tracking added between characters
                          (multiples of em). Default 0.18 spreads
                          chars enough to read as a label.
        direction:        bodyPr `vert` attribute. "vert270" = text
                          reads bottom-to-top (left edge), "vert" =
                          top-to-bottom (right edge), "horz" = no
                          rotation (debugging fallback).

    Returns:
        the textbox shape.

    Notes:
        * Place this zone with `role: "act-label"` in layoutSpec.zones
          so the title-choreography validator can recognize it.
        * Use the same role/text consistently across the deck — the
          point is structural rhythm, not per-slide labels.
    """
    from pptx.oxml.ns import qn

    if color is None:
        color = RGBColor(0x64, 0x74, 0x8B)

    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = False
    tf.auto_size = None

    # Apply vertical text direction at the bodyPr level.
    bodyPr = tf._txBody.find(qn("a:bodyPr"))
    if bodyPr is not None:
        bodyPr.set("vert", direction)
        bodyPr.set("wrap", "none")
        # Tight margins so the rotated text hugs the slide edge.
        bodyPr.set("lIns", "0")
        bodyPr.set("rIns", "0")
        bodyPr.set("tIns", "0")
        bodyPr.set("bIns", "0")

    rendered = (text or "").upper()
    # Letter-spacing via spc attribute (hundredths of a point).
    spc_val = int(letter_spacing_em * font_size_pt * 100)

    p = tf.paragraphs[0]
    p.text = rendered
    p.alignment = PP_ALIGN.LEFT
    run = p.runs[0]
    run.font.name = font_name
    run.font.size = Pt(font_size_pt)
    run.font.bold = bool(bold)
    run.font.color.rgb = color

    # Apply letter-spacing on the run's <a:rPr> via the spc attribute.
    rPr = run._r.find(qn("a:rPr"))
    if rPr is None:
        rPr = run._r.makeelement(qn("a:rPr"), {})
        run._r.insert(0, rPr)
    rPr.set("spc", str(spc_val))

    return txBox


def add_wide_tracking_title(slide, left, top, width, height, text,
                            color=None, font_name="Segoe UI",
                            font_size_pt=28, bold=True,
                            letter_spacing_em=0.20, uppercase=True,
                            alignment=PP_ALIGN.LEFT,
                            anchor=MSO_ANCHOR.TOP):
    """Add a title with extreme letter-spacing for a luxury/minimal feel.

    Implements the **wide-tracking** title move: extreme character
    spacing + uppercase transforms a standard heading into an
    editorial label that reads as designed, not templated. Best
    paired with `authoritative` or `inspirational` registers.

    Args:
        slide:            pptx slide object.
        left/top/width/height: EMU coordinates (use Inches()).
        text:             the title string. Uppercased if `uppercase` is True.
        color:            RGBColor; defaults to dark grey.
        font_name:        typeface; works best with geometric sans
                          (DM Sans, Inter, Montserrat).
        font_size_pt:     point size; 18–36pt typical for tracked titles.
        bold:             default True.
        letter_spacing_em: tracking (multiples of em). 0.15–0.30 is the
                          editorial sweet spot; >0.35 becomes hard to read.
        uppercase:        whether to force uppercase. Default True.
        alignment:        paragraph alignment.
        anchor:           vertical anchor.

    Returns:
        the textbox shape.

    Notes:
        * Place this zone with `tracking: "wide"` or
          `letterSpacing: <value>` in layoutSpec.zones so the
          title-choreography validator can recognize it.
    """
    from pptx.oxml.ns import qn

    if color is None:
        color = RGBColor(0x23, 0x23, 0x23)

    rendered = (text or "").upper() if uppercase else (text or "")

    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    try:
        tf.vertical_anchor = anchor
    except Exception:
        pass

    p = tf.paragraphs[0]
    p.text = rendered
    p.alignment = alignment

    run = p.runs[0]
    run.font.name = font_name
    run.font.size = Pt(font_size_pt)
    run.font.bold = bool(bold)
    run.font.color.rgb = color

    # Apply letter-spacing via the spc attribute (hundredths of a point).
    spc_val = int(letter_spacing_em * font_size_pt * 100)
    rPr = run._r.find(qn("a:rPr"))
    if rPr is None:
        rPr = run._r.makeelement(qn("a:rPr"), {})
        run._r.insert(0, rPr)
    rPr.set("spc", str(spc_val))

    return txBox


def add_drop_cap(slide, left, top, cap_width, cap_height,
                 char, rest_text,
                 cap_color=None, cap_font_name="Segoe UI",
                 cap_font_size_pt=72, cap_bold=True,
                 body_left=None, body_top=None,
                 body_width=None, body_height=None,
                 body_font_name="Segoe UI", body_font_size_pt=16,
                 body_color=None, body_alignment=PP_ALIGN.LEFT):
    """Add a drop-cap initial letter with adjacent body text.

    Implements the **drop-cap** title move: the first letter is
    rendered at 2–4× body size in a separate textbox, creating a
    magazine editorial feel. The remaining text flows beside it.

    Args:
        slide:           pptx slide object.
        left/top:        EMU coordinates of the cap letter box.
        cap_width/cap_height: EMU size of the cap letter box.
        char:            the initial character (typically 1 char).
        rest_text:       the remaining title/body text after the cap.
        cap_color:       RGBColor for the cap; defaults to dark grey.
        cap_font_name:   typeface for the cap.
        cap_font_size_pt: point size for the cap (48–96pt typical).
        cap_bold:        default True.
        body_left/top/width/height: EMU coordinates for the body text
                         box. If None, positioned to the right of cap
                         with a small gutter.
        body_font_name:  typeface for the body text.
        body_font_size_pt: point size for the body.
        body_color:      RGBColor for the body; defaults to dark grey.
        body_alignment:  paragraph alignment for the body.

    Returns:
        tuple of (cap_shape, body_shape).

    Notes:
        * Place the cap zone with `role` containing `drop-cap` in
          layoutSpec.zones so the validator can recognize it.
    """
    if cap_color is None:
        cap_color = RGBColor(0x23, 0x23, 0x23)
    if body_color is None:
        body_color = RGBColor(0x33, 0x33, 0x33)

    # Cap letter box.
    cap_box = slide.shapes.add_textbox(left, top, cap_width, cap_height)
    tf_cap = cap_box.text_frame
    tf_cap.word_wrap = False
    tf_cap.auto_size = None
    tf_cap.margin_left = 0
    tf_cap.margin_right = 0
    tf_cap.margin_top = 0
    tf_cap.margin_bottom = 0

    p_cap = tf_cap.paragraphs[0]
    p_cap.text = char
    p_cap.alignment = PP_ALIGN.CENTER
    run_cap = p_cap.runs[0]
    run_cap.font.name = cap_font_name
    run_cap.font.size = Pt(cap_font_size_pt)
    run_cap.font.bold = bool(cap_bold)
    run_cap.font.color.rgb = cap_color
    p_cap.line_spacing = 1.0

    # Body text box — default position to the right of the cap.
    gutter = Inches(0.15)
    if body_left is None:
        body_left = left + cap_width + gutter
    if body_top is None:
        body_top = top
    if body_width is None:
        body_width = Inches(8.0)  # caller should override
    if body_height is None:
        body_height = cap_height

    body_box = slide.shapes.add_textbox(body_left, body_top,
                                        body_width, body_height)
    tf_body = body_box.text_frame
    tf_body.word_wrap = True
    tf_body.auto_size = None

    p_body = tf_body.paragraphs[0]
    p_body.text = rest_text or ""
    p_body.alignment = body_alignment
    run_body = p_body.runs[0] if p_body.runs else p_body.add_run()
    run_body.font.name = body_font_name
    run_body.font.size = Pt(body_font_size_pt)
    run_body.font.color.rgb = body_color

    return cap_box, body_box


def add_split_highlight(slide, text_left, text_top, text_width, text_height,
                        full_text, highlight_words,
                        text_font_name="Segoe UI", text_font_size_pt=28,
                        text_color=None, highlight_fill=None,
                        highlight_height_factor=0.85,
                        highlight_y_offset_factor=0.15,
                        alignment=PP_ALIGN.LEFT,
                        anchor=MSO_ANCHOR.TOP):
    """Add a title with colored rectangle highlights behind key words.

    Implements the **split-highlight** title move: 1–2 key words in
    the title are backed by a colored rectangle (highlighter marker
    effect), creating a strong focal point. The highlight shapes are
    rendered BEFORE the text box so they appear behind the text.

    Args:
        slide:           pptx slide object.
        text_left/top/width/height: EMU coordinates for the title text box.
        full_text:       the complete title string.
        highlight_words: list of 1–2 exact substrings to highlight.
        text_font_name:  typeface.
        text_font_size_pt: point size.
        text_color:      RGBColor for the text; defaults to dark grey.
        highlight_fill:  RGBColor for the highlight rectangles;
                         defaults to a warm accent yellow.
        highlight_height_factor: fraction of line height for the
                         highlight rectangle (0.85 = 85% of line).
        highlight_y_offset_factor: vertical offset as fraction of
                         line height (pushes highlight toward baseline).
        alignment:       paragraph alignment.
        anchor:          vertical anchor.

    Returns:
        tuple of (text_shape, list_of_highlight_shapes).

    Notes:
        * Place a shape zone with `role` containing `highlight` in
          layoutSpec.zones so the validator can recognize this move.
        * The highlight rectangles are approximate — they use
          character-count proportional positioning. For precise
          positioning, measure the rendered text in Step 8 QA.
        * Highlight is best with 1–2 short words; longer phrases
          risk overflow.
    """
    if text_color is None:
        text_color = RGBColor(0x23, 0x23, 0x23)
    if highlight_fill is None:
        highlight_fill = RGBColor(0xFF, 0xE0, 0x82)  # warm accent yellow

    highlight_shapes = []
    line_height_emu = Pt(text_font_size_pt * 1.3)
    char_width_emu = Pt(text_font_size_pt * 0.55)  # approximate average

    for word in (highlight_words or []):
        idx = full_text.find(word)
        if idx < 0:
            continue
        # Approximate horizontal position based on character count.
        hl_left = text_left + int(idx * char_width_emu)
        hl_width = int(len(word) * char_width_emu)
        hl_height = int(line_height_emu * highlight_height_factor)
        hl_top = text_top + int(line_height_emu * highlight_y_offset_factor)

        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, hl_left, hl_top, hl_width, hl_height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = highlight_fill
        shape.line.fill.background()  # no border
        highlight_shapes.append(shape)

    # Text box on top of highlights.
    txBox = slide.shapes.add_textbox(text_left, text_top,
                                     text_width, text_height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    try:
        tf.vertical_anchor = anchor
    except Exception:
        pass

    p = tf.paragraphs[0]
    p.text = full_text
    p.alignment = alignment
    run = p.runs[0]
    run.font.name = text_font_name
    run.font.size = Pt(text_font_size_pt)
    run.font.color.rgb = text_color

    return txBox, highlight_shapes


# ── Accent motif (Phase B — see visual-tone-mapping.md § Accent imagery) ──

def add_accent_motif(slide, img_path, left, top, width, height,
                     *, fallback_fill=None, transparency=None,
                     placeholder_text=None, fit="stretch"):
    """Place a small decorative image (≤ 30% of canvas) without scrim
    or clip. Use for accent imagery: abstract textures, color washes,
    organic motifs, gradient bands placed at the slide periphery to
    add visual interest without competing with the content.

    This helper differs from `add_image_safe` in three ways suited
    to the accent role:

      1. The placeholder fallback is a soft tinted block (not a grey
         rectangle), so a missing accent image degrades to "subtle
         color shape" rather than "obvious missing-asset placeholder".
      2. No filename is rendered into the placeholder \u2014 accent
         motifs should never reveal asset metadata.
      3. Optional `transparency` is applied so the motif sits
         visually behind / beside content without grabbing attention.

    Args:
        slide:            pptx slide object.
        img_path:         path to the motif image (str / Path) or None.
        left/top/width/height: EMU coordinates (use Inches()).
        fallback_fill:    RGBColor used when image is missing. Defaults
                          to a soft warm grey (E8EBF0).
        transparency:     0\u2013100 percent. None or 0 = opaque. For
                          atmospheric edge motifs, 25\u201340 typically
                          reads as "ambient" without disappearing.
        placeholder_text: optional caption rendered ONLY in the missing-
                          image fallback path (debugging aid; defaults
                          to None so the placeholder stays clean).
        fit:              "stretch" (default) fills the zone exactly,
                          ignoring the image aspect ratio. Recommended
                          for abstract washes / textures / color fields
                          where there is no subject to distort.
                          "contain" preserves aspect, centering the
                          image inside the zone (use for accent motifs
                          that have a recognizable composition).

    Returns:
        the image / placeholder shape, or None if width/height was 0.
    """
    if width == 0 or height == 0:
        return None

    if fallback_fill is None:
        fallback_fill = RGBColor(0xE8, 0xEB, 0xF0)

    img_path_obj = (Path(img_path) if img_path and not isinstance(img_path, Path)
                    else img_path)

    # Missing-asset fallback: soft tinted block.
    if not img_path_obj or not Path(img_path_obj).exists():
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = fallback_fill
        shape.line.fill.background()
        if transparency:
            _apply_shape_alpha(shape, 1.0 - max(0.0, min(100.0, transparency)) / 100.0)
        if placeholder_text:
            add_textbox(slide, left + Inches(0.1), top + Inches(0.1),
                        width - Inches(0.2), Inches(0.3),
                        placeholder_text, font_size=9,
                        color=RGBColor(0x99, 0x99, 0x99))
        return shape

    # Real-image path. Default fit is "stretch" — for abstract
    # washes / textures the zone aspect should win. "contain"
    # preserves the image aspect (centered) when the motif itself
    # has a composition that must not be distorted.
    if fit == "contain":
        try:
            from PIL import Image
            img = Image.open(img_path_obj)
            img_w, img_h = img.size
            aspect = img_w / img_h
            zone_aspect = width / height
            if aspect > zone_aspect:
                display_w = width
                display_h = int(width / aspect)
            else:
                display_h = height
                display_w = int(height * aspect)
            offset_x = (width - display_w) // 2
            offset_y = (height - display_h) // 2
            pic = slide.shapes.add_picture(
                str(img_path_obj),
                left + offset_x, top + offset_y,
                display_w, display_h,
            )
        except ImportError:
            pic = slide.shapes.add_picture(
                str(img_path_obj), left, top, width, height)
    else:
        # stretch — fill the zone exactly.
        pic = slide.shapes.add_picture(
            str(img_path_obj), left, top, width, height)

    if transparency:
        # Transparency on a Picture requires direct XML on the
        # blipFill alphaModFix. Apply via low-level helper.
        from pptx.oxml.ns import qn
        blip = pic._element.find(".//" + qn("a:blip"))
        if blip is not None:
            from lxml import etree
            alpha_pct = max(0.0, min(100.0, transparency))
            alpha_val = str(int((100 - alpha_pct) * 1000))
            # alphaModFix expects amt = remaining opacity in 1/1000 of percent
            existing = blip.find(qn("a:alphaModFix"))
            if existing is not None:
                blip.remove(existing)
            am = etree.SubElement(blip, qn("a:alphaModFix"))
            am.set("amt", alpha_val)

    return pic


# ── Formula (mathtext → PNG) ─────────────────────────────────────
#
# Phase 1 of the formula-zone proposal. See
# `sessions/formula-svg-design-2026-05-01.md` for the full design.
#
# **Phase 1 artifact = PNG @ 300dpi.** python-pptx insertion goes
# through Pillow, which does not understand SVG. Native SVG embedding
# in .pptx requires direct OOXML manipulation (`<asvg:svgBlip>`
# extension on the `<a:blip>`) plus a raster PNG fallback for
# compatibility — that is tracked as Phase 1.5 follow-up. PNG @ 300dpi
# in standard 16:9 zone sizes (≥ 0.4 in tall) is visually crisp at
# every zoom level a presenter actually uses.
#
# `render_formula_svg(...)` keeps SVG output capability for
# offline/archival use and the future XML-injection path; the
# caller chooses by file suffix on `out_path`.
#
# Three responsibilities:
#   * `validate_mathtext(source)`  — pure, no I/O. Used by s06 to
#     dry-run-parse every formulaSource at validation time.
#   * `render_formula_svg(...)`    — render one mathtext expression
#     (writes SVG or PNG depending on out_path suffix). No slide.
#   * `add_formula(...)`           — compose: render-with-cache as PNG,
#     fit into the zone, insert as a picture; fall back to a textbox
#     when render fails.
#
# All three live here so per-session `s07-build.py` scripts only need
# the existing single import from this module.

VALID_MATHTEXT_FONTSETS = (
    "dejavusans", "dejavuserif", "cm", "stix", "stixsans",
)
# Default to "cm" (Computer Modern) — matplotlib mathtext's LaTeX-faithful
# fontset. Math reads as classical LaTeX typesetting (italic Latin, upright
# Greek, math-italic glyphs) — the convention every paper-trained reader
# already maps onto. Override per-deck via
# s05b-style-policy.json → typography.mathFontset when a sans/serif math
# face suits the deck's typographic register better.
DEFAULT_MATHTEXT_FONTSET = "cm"

_MATHTEXT_PARSER_CACHE = {}


def validate_mathtext(source):
    """Dry-run-parse a mathtext expression body (no `$…$` delimiters).

    Returns:
        (ok: bool, err: str|None). On failure, `err` carries the
        matplotlib parser message — surface it directly to the author.

    Notes:
        * Forbidden inputs (delimiters around the expression) are NOT
          rejected here — that is `s06_validate_content.py`'s job
          (different error voice). This function is the syntactic
          mathtext check only.
        * The parser is cached per-process; cost ≈ 1ms after warmup.
    """
    if not isinstance(source, str) or not source.strip():
        return False, "formula source is empty"
    try:
        from matplotlib import mathtext as _mt
        parser = _MATHTEXT_PARSER_CACHE.get("path")
        if parser is None:
            parser = _mt.MathTextParser("path")
            _MATHTEXT_PARSER_CACHE["path"] = parser
        # MathTextParser.parse expects the wrapped form.
        parser.parse(f"${source}$")
        return True, None
    except ImportError as e:
        return False, f"matplotlib unavailable: {e}"
    except Exception as e:
        # mathtext raises ValueError, plus its own ParseFatalException.
        return False, str(e)


def _normalize_color_hex(value):
    """Coerce hex string / RGBColor / (r,g,b) tuple → '#RRGGBB' string."""
    if isinstance(value, RGBColor):
        return "#{:02X}{:02X}{:02X}".format(value[0], value[1], value[2])
    if isinstance(value, (list, tuple)) and len(value) == 3:
        return "#{:02X}{:02X}{:02X}".format(*value)
    if isinstance(value, str):
        v = value.lstrip("#")
        if len(v) != 6:
            raise ValueError(f"hex color must be 6 chars: {value!r}")
        return "#" + v.upper()
    raise ValueError(f"cannot normalize color: {value!r}")


def render_formula_svg(source, *, color_hex, font_pt, out_path,
                       fontset=DEFAULT_MATHTEXT_FONTSET,
                       png_fallback_dpi=300):
    """Render a mathtext expression body to SVG (or PNG on fallback).

    Args:
        source:           mathtext expression body, NO `$…$` delimiters.
        color_hex:        '#RRGGBB' string (or any form accepted by
                          `_normalize_color_hex`). Baked into the SVG.
        font_pt:          font size in points.
        out_path:         destination path. Suffix decides format
                          (`.svg` primary, `.png` is the fallback path).
        fontset:          one of `VALID_MATHTEXT_FONTSETS`.
        png_fallback_dpi: only used when out_path ends with `.png`.

    Returns:
        Path to the written file (matches `out_path`).

    Raises:
        ValueError on a parser failure (after `validate_mathtext` should
        have caught it at s06; raising here means a bug or the author
        bypassed the validator).
    """
    if fontset not in VALID_MATHTEXT_FONTSETS:
        raise ValueError(
            f"fontset {fontset!r} not in {VALID_MATHTEXT_FONTSETS}"
        )

    color = _normalize_color_hex(color_hex)
    out_path = Path(out_path)
    out_path.parent.mkdir(parents=True, exist_ok=True)

    import matplotlib
    matplotlib.use("Agg")
    import matplotlib.pyplot as plt
    from matplotlib import rcParams

    fmt = "png" if out_path.suffix.lower() == ".png" else "svg"

    rcParams["svg.fonttype"] = "path"     # embed glyphs as paths
    rcParams["mathtext.fontset"] = fontset

    fig = plt.figure(figsize=(0.01, 0.01))
    try:
        fig.text(
            0, 0,
            f"${source}$",
            fontsize=font_pt,
            color=color,
        )
        save_kwargs = dict(
            format=fmt,
            bbox_inches="tight",
            pad_inches=0.02,
            transparent=True,
        )
        if fmt == "png":
            save_kwargs["dpi"] = png_fallback_dpi
        fig.savefig(str(out_path), **save_kwargs)
    finally:
        plt.close(fig)

    return out_path


def _measure_svg_extent_emu(svg_path):
    """Read width/height from a matplotlib-saved SVG, return EMU.

    matplotlib writes `<svg width="…pt" height="…pt" viewBox="…">`.
    1pt = 12700 EMU. Returns (w_emu, h_emu) as ints.

    Falls back to `(None, None)` on any parse failure — caller should
    treat as "unable to measure" and skip the fit-down branch.
    """
    try:
        import re
        text = Path(svg_path).read_text(encoding="utf-8", errors="ignore")
        m = re.search(
            r'<svg[^>]*\bwidth="([\d.]+)pt"[^>]*\bheight="([\d.]+)pt"',
            text,
        )
        if not m:
            return None, None
        w_pt = float(m.group(1))
        h_pt = float(m.group(2))
        return int(w_pt * 12700), int(h_pt * 12700)
    except Exception:
        return None, None


def _measure_image_extent_emu(img_path, *, png_dpi=300):
    """Return (w_emu, h_emu) for a rendered formula artifact.

    Dispatches by suffix: SVG → parse `width="…pt"`. PNG → Pillow
    read px dimensions, divide by dpi, multiply by 914400 EMU/in.
    Returns (None, None) when measurement is unavailable.
    """
    p = Path(img_path)
    suffix = p.suffix.lower()
    if suffix == ".svg":
        return _measure_svg_extent_emu(p)
    if suffix == ".png":
        try:
            from PIL import Image
            with Image.open(p) as img:
                w_px, h_px = img.size
                # Prefer the file's embedded dpi over the caller hint.
                dpi = img.info.get("dpi")
                if dpi and isinstance(dpi, (tuple, list)) and dpi[0] > 0:
                    dpi_x, dpi_y = float(dpi[0]), float(dpi[1] or dpi[0])
                else:
                    dpi_x = dpi_y = float(png_dpi)
                w_emu = int(w_px / dpi_x * 914400)
                h_emu = int(h_px / dpi_y * 914400)
                return w_emu, h_emu
        except Exception:
            return None, None
    return None, None


def _set_picture_alt(pic, alt_text):
    """Write `alt_text` to the picture shape's <p:nvPicPr>/<p:cNvPr>.descr."""
    if not alt_text:
        return
    try:
        nv = pic._element.find(
            "./{http://schemas.openxmlformats.org/presentationml/2006/main}nvPicPr"
        )
        if nv is None:
            return
        cnv = nv.find(
            "./{http://schemas.openxmlformats.org/presentationml/2006/main}cNvPr"
        )
        if cnv is not None:
            cnv.set("descr", alt_text)
    except Exception:
        # Alt text is best-effort; never let a11y stamping break the build.
        pass


# PowerPoint's documented marker for the SVG blip extension on a picture.
_SVG_EXT_URI = "{96DAC541-7B7A-43D3-8B79-37D633B846F1}"
_ASVG_NS = "http://schemas.microsoft.com/office/drawing/2016/SVG/main"
_R_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"


def _register_svg_image_part(slide, svg_path):
    """Register an SVG file as an image Part on the package and relate it to the slide.

    Returns the relationship id (rId) the slide uses to reference the part,
    or None on any failure (caller should treat as "no SVG attached" and
    fall through to PNG-only).
    """
    try:
        from pptx.opc.constants import RELATIONSHIP_TYPE as RT
        from pptx.opc.package import Part

        pkg = slide.part.package
        partname = pkg.next_partname("/ppt/media/image%d.svg")
        with open(svg_path, "rb") as f:
            blob = f.read()
        svg_part = Part(partname, "image/svg+xml", pkg, blob)
        return slide.part.relate_to(svg_part, RT.IMAGE)
    except Exception:
        return None


def _inject_svg_blip_extension(pic, svg_rid):
    """Append an `<asvg:svgBlip>` extension to the picture's `<a:blip>`.

    Mirrors the OOXML pattern PowerPoint emits when inserting a vector SVG:

        <a:blip r:embed="<png-rId>">
          <a:extLst>
            <a:ext uri="{96DAC541-...}">
              <asvg:svgBlip xmlns:asvg="..." r:embed="<svg-rId>"/>
            </a:ext>
          </a:extLst>
        </a:blip>

    PowerPoint 2019+ renders the vector; older viewers fall back to the
    PNG referenced by the parent blip's `r:embed`. Returns True on success.
    """
    try:
        from lxml import etree

        from pptx.oxml.ns import qn

        blip = pic._element.find(".//" + qn("a:blip"))
        if blip is None:
            return False

        ext_lst = blip.find(qn("a:extLst"))
        if ext_lst is None:
            ext_lst = etree.SubElement(blip, qn("a:extLst"))

        ext = etree.SubElement(ext_lst, qn("a:ext"))
        ext.set("uri", _SVG_EXT_URI)

        svg_blip = etree.SubElement(
            ext, "{%s}svgBlip" % _ASVG_NS, nsmap={"asvg": _ASVG_NS}
        )
        svg_blip.set("{%s}embed" % _R_NS, svg_rid)
        return True
    except Exception:
        return False


def add_formula(slide, source, left, top, width, height, *,
                color="222222",
                font_pt=24,
                fontset=DEFAULT_MATHTEXT_FONTSET,
                alignment="center",
                vertical_align="middle",
                fit_mode="scale-to-zone",
                natural_promote_threshold=0.40,
                cache_path=None,
                fallback_text=None,
                alt=None,
                png_fallback_dpi=300,
                build_log=None):
    """Render `source` as SVG and place it inside the zone.

    See `sessions/formula-svg-design-2026-05-01.md` for the full
    contract. Mirrors `add_image_safe` in shape: same (slide, l, t, w, h)
    preamble, missing-input fallback to a placeholder, never raises on
    a bad render.

    Failure ladder:
        1. `validate_mathtext` fails  → write `fallback_text` as a
           textbox in the zone and log.
        2. SVG render fails           → retry as PNG@`png_fallback_dpi`.
        3. PNG render also fails      → fallback_text textbox.

    `cache_path` (Path) is the deterministic on-disk artifact location
    (e.g. `sessions/<sid>/formulas/<taskId>-<role>.svg`). When the file
    already exists, the helper skips re-rendering — callers should
    delete cached files when the source / color / font / fontset
    change (the deck-level `render_formulas()` pre-pass owns this).

    `build_log`, if provided, is a list mutated in place with one dict
    per call recording `{taskId, role, action, reason}` so the caller
    can summarize for the s07 build report.

    Returns the inserted picture shape (or fallback textbox shape).
    """
    log_entry = {
        "role": None,           # caller can stamp this after return
        "action": None,         # 'svg' | 'svg-cached' | 'png-fallback' | 'text-fallback'
        "fit_mode_applied": fit_mode,
        "reason": None,
    }

    def _emit_fallback(reason):
        log_entry["action"] = "text-fallback"
        log_entry["reason"] = reason
        text = fallback_text or source or ""
        return add_textbox(
            slide, left, top, width, height,
            text,
            font_size=max(10, int(font_pt * 0.6)),
            color=resolve_color(color)
            if isinstance(color, str) and len(color.lstrip("#")) == 6
            else RGBColor(0x22, 0x22, 0x22),
            alignment={
                "left": PP_ALIGN.LEFT,
                "right": PP_ALIGN.RIGHT,
                "center": PP_ALIGN.CENTER,
            }.get(alignment, PP_ALIGN.CENTER),
            anchor={
                "top": MSO_ANCHOR.TOP,
                "bottom": MSO_ANCHOR.BOTTOM,
                "middle": MSO_ANCHOR.MIDDLE,
            }.get(vertical_align, MSO_ANCHOR.MIDDLE),
        )

    # 1. Syntactic check (cheap).
    ok, parse_err = validate_mathtext(source or "")
    if not ok:
        shape = _emit_fallback(f"mathtext parse error: {parse_err}")
        if build_log is not None:
            build_log.append(log_entry)
        return shape

    color_hex = _normalize_color_hex(color)

    # 2. Render (or reuse cache).
    artifact = None
    used_cache = False
    if cache_path:
        cache_path = Path(cache_path)
        if cache_path.exists() and cache_path.stat().st_size > 0:
            artifact = cache_path
            used_cache = True
            log_entry["action"] = "svg-cached" if cache_path.suffix.lower() == ".svg" else "png-cached"

    if artifact is None:
        try:
            target = Path(cache_path) if cache_path else Path(f"_formula_{id(source)}.png")
            # Phase 1 primary artifact = PNG (python-pptx cannot insert SVG).
            if target.suffix.lower() not in (".png", ".svg"):
                target = target.with_suffix(".png")
            artifact = render_formula_svg(
                source,
                color_hex=color_hex,
                font_pt=font_pt,
                out_path=target,
                fontset=fontset,
                png_fallback_dpi=png_fallback_dpi,
            )
            log_entry["action"] = "png" if target.suffix.lower() == ".png" else "svg"
        except Exception as render_err:
            shape = _emit_fallback(f"render failed: {render_err}")
            if build_log is not None:
                build_log.append(log_entry)
            return shape

    # 3. Measure → fit-mode resolution.
    nat_w, nat_h = _measure_image_extent_emu(
        artifact, png_dpi=png_fallback_dpi,
    )
    effective_fit = fit_mode

    if effective_fit == "natural" and nat_w and nat_h:
        bind_ratio = min(nat_w / width, nat_h / height)
        if bind_ratio < natural_promote_threshold:
            effective_fit = "scale-to-zone"
            log_entry["fit_mode_applied"] = "scale-to-zone(auto-promoted)"
            log_entry["reason"] = (
                f"natural extent {bind_ratio:.0%} < {natural_promote_threshold:.0%} "
                f"of zone — auto-promoted to scale-to-zone"
            )

    # 4. Compute display rect, preserving aspect.
    if effective_fit == "scale-to-zone" and nat_w and nat_h:
        nat_aspect = nat_w / nat_h
        zone_aspect = width / height
        if nat_aspect > zone_aspect:
            disp_w = width
            disp_h = int(width / nat_aspect)
        else:
            disp_h = height
            disp_w = int(height * nat_aspect)
    elif effective_fit == "natural" and nat_w and nat_h:
        disp_w, disp_h = nat_w, nat_h
        # Clamp if natural exceeds the zone (should not happen but be safe).
        if disp_w > width or disp_h > height:
            disp_w, disp_h = width, height
    else:
        # Could not measure — fit to zone exactly (PNG path or parse-failed SVG).
        disp_w, disp_h = width, height

    # 5. Position by alignment / verticalAlign within zone.
    if alignment == "left":
        offset_x = 0
    elif alignment == "right":
        offset_x = width - disp_w
    else:
        offset_x = (width - disp_w) // 2

    if vertical_align == "top":
        offset_y = 0
    elif vertical_align == "bottom":
        offset_y = height - disp_h
    else:
        offset_y = (height - disp_h) // 2

    # 6. Insert.
    pic = slide.shapes.add_picture(
        str(artifact),
        left + offset_x, top + offset_y,
        disp_w, disp_h,
    )
    _set_picture_alt(pic, alt)

    # 6b. Phase 1.5 — attach SVG sidecar as a vector overlay so PowerPoint
    # 2019+ renders the vector while older viewers see the inserted PNG.
    # Sidecar is the same basename with .svg extension. Failure is silent
    # and non-fatal — the PNG insertion above is always the baseline.
    artifact_path = Path(artifact)
    if artifact_path.suffix.lower() == ".png":
        svg_sidecar = artifact_path.with_suffix(".svg")
        if svg_sidecar.exists() and svg_sidecar.stat().st_size > 0:
            svg_rid = _register_svg_image_part(slide, svg_sidecar)
            if svg_rid and _inject_svg_blip_extension(pic, svg_rid):
                if log_entry["action"] in ("png", "png-cached"):
                    log_entry["action"] = log_entry["action"] + "+svg"

    if build_log is not None:
        build_log.append(log_entry)
    return pic


def resolve_formula_color(style_policy, color_token=None):
    """Resolve the formula color per the design's default chain.

    Priority (per sessions/formula-svg-design-2026-05-01.md § 2.1):
        1. explicit `color_token` from the s05 zone (any palette token)
        2. style_policy.palette.body
        3. style_policy.palette.ink
        4. '#222222' (hard fallback)

    Returns a 6-char hex string ('#RRGGBB' uppercase, with leading '#').
    """
    palette = ((style_policy or {}).get("palette") or {})
    if color_token:
        token_val = palette.get(color_token)
        if token_val:
            return _normalize_color_hex(token_val)
    for fallback_key in ("body", "ink"):
        token_val = palette.get(fallback_key)
        if token_val:
            return _normalize_color_hex(token_val)
    return "#222222"


def resolve_formula_fontset(style_policy):
    """Read style-policy → typography.mathFontset, default 'dejavusans'."""
    typo = ((style_policy or {}).get("typography") or {})
    fs = typo.get("mathFontset", DEFAULT_MATHTEXT_FONTSET)
    if fs not in VALID_MATHTEXT_FONTSETS:
        return DEFAULT_MATHTEXT_FONTSET
    return fs


def render_formulas(plan, content, *, session_dir, style_policy=None,
                    default_font_pt=24, png_dpi=300, build_log=None):
    """Pre-pass: render every formula slot in `content` to disk.

    Walks `s05-slide-visual-design.json` for `formula` zones, finds
    the matching slot in `s06-slide-content.json`, resolves color +
    fontset from `style_policy`, renders both a vector
    `<session_dir>/formulas/<taskId>-<role>.svg` (Phase 1.5 archival
    + on-slide vector overlay) and a raster
    `<session_dir>/formulas/<taskId>-<role>.png` (raster baseline that
    `add_picture` can actually insert), returning a dict keyed by
    `(taskId, role) → png_cache_path` so per-slide builders can pass
    it directly to `add_formula(cache_path=...)`. `add_formula` then
    finds the SVG sidecar by suffix swap and attaches it via the
    `<asvg:svgBlip>` extension.

    Cache invalidation is automatic via a sidecar `.sig` file holding
    `source|color|font_pt|fontset|dpi`. Any drift triggers a re-render
    of both files. SVG render failure is non-fatal — the PNG still
    powers the slide and a `svg-render-failed` entry is logged.

    `build_log` is mutated in place with one entry per render.
    """
    cache_map = {}
    formulas_dir = Path(session_dir) / "formulas"
    formulas_dir.mkdir(parents=True, exist_ok=True)

    fontset = resolve_formula_fontset(style_policy)

    plan_slides = {s.get("taskId"): s for s in plan.get("slides", [])}
    for cs in content.get("slides", []):
        tid = cs.get("taskId")
        plan_slide = plan_slides.get(tid)
        if not plan_slide:
            continue
        zones_by_role = {
            z.get("role"): z
            for z in plan_slide.get("layoutSpec", {}).get("zones", [])
            if z.get("type") == "formula"
        }
        if not zones_by_role:
            continue
        for slot in cs.get("slots", []):
            role = slot.get("zone")
            zone = zones_by_role.get(role)
            if not zone:
                continue
            source = slot.get("formulaSource") or ""
            if not source:
                continue
            color_hex = resolve_formula_color(
                style_policy, zone.get("colorToken")
            )
            font_pt = float(zone.get("fontPt") or default_font_pt)
            png_out = formulas_dir / f"{tid}-{role}.png"
            svg_out = formulas_dir / f"{tid}-{role}.svg"

            sig = f"{source}|{color_hex}|{font_pt}|{fontset}|dpi={png_dpi}"
            sig_file = png_out.with_suffix(".png.sig")
            if png_out.exists() and svg_out.exists() and sig_file.exists():
                try:
                    if sig_file.read_text(encoding="utf-8") == sig:
                        cache_map[(tid, role)] = png_out
                        if build_log is not None:
                            build_log.append({
                                "taskId": tid, "role": role,
                                "action": "png+svg-cached",
                                "path": str(png_out),
                            })
                        continue
                except Exception:
                    pass
                # signature drift — re-render below

            png_ok = False
            svg_ok = False
            try:
                render_formula_svg(
                    source,
                    color_hex=color_hex,
                    font_pt=font_pt,
                    out_path=png_out,
                    fontset=fontset,
                    png_fallback_dpi=png_dpi,
                )
                png_ok = True
            except Exception as err:
                if build_log is not None:
                    build_log.append({
                        "taskId": tid, "role": role,
                        "action": "png-render-failed",
                        "reason": str(err),
                    })
                # No PNG = no insertion path — skip SVG too.
                continue

            try:
                render_formula_svg(
                    source,
                    color_hex=color_hex,
                    font_pt=font_pt,
                    out_path=svg_out,
                    fontset=fontset,
                    png_fallback_dpi=png_dpi,
                )
                svg_ok = True
            except Exception as err:
                if build_log is not None:
                    build_log.append({
                        "taskId": tid, "role": role,
                        "action": "svg-render-failed",
                        "reason": str(err),
                    })

            sig_file.write_text(sig, encoding="utf-8")
            cache_map[(tid, role)] = png_out
            if build_log is not None:
                action = "png+svg" if (png_ok and svg_ok) else "png"
                build_log.append({
                    "taskId": tid, "role": role,
                    "action": action, "path": str(png_out),
                })

    return cache_map
