#!/usr/bin/env python3
"""Step 8 per-slide QA — precompute zone metrics.

Reads the built ``output.pptx`` and produces ``s08-zone-metrics.json``
containing **raw measurements only**. The VLM detection subagent
consumes these metrics together with rendered JPEGs and the §3.1
dimension grading criteria from workflow-step8.md to assign grades.

Strict design rule: this script outputs numbers and factual
observations, never pass/fail judgments. There are no boolean
fields like ``is_overflow`` or ``distorted``. Examples:

* ``wcagContrast.fgVsBg = 3.8`` ← raw ratio, NOT "LOW_CONTRAST: true"
* ``image.aspectDelta = 0.41``  ← raw delta,  NOT "distorted: true"
* ``outOfCanvas: [{zoneId, overhangIn}]`` ← list of facts, NOT severity

Schema: ``$COMPOSER_SKILL/schemas/s08-zone-metrics.schema.json``.

Usage:
  python s08_precompute_metrics.py <session_dir>
  python s08_precompute_metrics.py .                       # from inside session
  python s08_precompute_metrics.py . --output build2.pptx  # non-default pptx
"""

from __future__ import annotations

import argparse
import datetime as _dt
import json
import re
import sys
from pathlib import Path
from typing import Any

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Emu

# Force UTF-8 stdout on Windows.
if sys.platform == "win32":
    try:
        sys.stdout.reconfigure(encoding="utf-8")
        sys.stderr.reconfigure(encoding="utf-8")
    except (AttributeError, OSError):
        pass


# ── Constants ────────────────────────────────────────────────────────

EMU_PER_INCH = 914400.0

PLACEHOLDER_TOKENS = [
    "xxxx",
    "lorem",
    "ipsum",
    "click to",
    "insert",
    "TBD",
    "placeholder",
    "sample text",
    "[placeholder]",
    "to be filled",
]

CODE_ZONE_NAME_HINTS = ("code", "term", "console", "snippet", "shell")

DEFAULT_BG = "#FFFFFF"

# ── v2 line-wrap estimator constants (Phase 1, 2026-05-02) ───────────

# Average glyph width as a fraction of fontPt (em). Coarse estimate;
# does not need precise font metrics — VLM is the final judge. Heading
# uses 0.65 (calibrated against bold display sans-serif digits "0"/"2"
# at 280 pt in a 5.0 in zone, which empirically wraps); body uses 0.50
# (proportional sans-serif average); mono is fixed-width.
CHAR_WIDTH_EM_HEADING = 0.65
CHAR_WIDTH_EM_BODY = 0.50
CHAR_WIDTH_EM_MONO = 0.60

# Default zone padding subtracted from each side when computing usable width.
DEFAULT_TEXT_INSET_IN = 0.1

# fontPt threshold for promoting body → heading classification.
HEADING_FONT_PT_THRESHOLD = 24.0

# Substrings in roleHint that promote a zone to "heading" classification
# regardless of fontPt. Lowercased on compare.
HEADING_ROLE_HINT_TOKENS = ("title", "numeral", "heading", "headline", "act-")


# ── Helpers ──────────────────────────────────────────────────────────


def emu_to_in(v: int | None) -> float:
    if v is None:
        return 0.0
    return round(v / EMU_PER_INCH, 4)


def rgb_to_hex(rgb: RGBColor | None) -> str | None:
    if rgb is None:
        return None
    try:
        return "#" + str(rgb).upper()
    except Exception:
        return None


def hex_to_rgb_tuple(h: str) -> tuple[int, int, int]:
    h = h.lstrip("#")
    return int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)


def relative_luminance(rgb: tuple[int, int, int]) -> float:
    """WCAG 2.x relative luminance of an sRGB color."""

    def chan(c: int) -> float:
        s = c / 255.0
        return s / 12.92 if s <= 0.03928 else ((s + 0.055) / 1.055) ** 2.4

    r, g, b = (chan(c) for c in rgb)
    return 0.2126 * r + 0.7152 * g + 0.0722 * b


def wcag_contrast(fg_hex: str, bg_hex: str) -> float:
    l1 = relative_luminance(hex_to_rgb_tuple(fg_hex))
    l2 = relative_luminance(hex_to_rgb_tuple(bg_hex))
    a, b = max(l1, l2), min(l1, l2)
    return round((a + 0.05) / (b + 0.05), 2)


# ── Color resolution (best-effort, returns None on inheritance gaps) ─


def _solid_fill_hex(fill_format) -> tuple[str | None, str | None]:
    """Return (hex_color, fill_type) for a shape/text fill format.

    Returns (None, fill_type) when the fill is solid but uses a theme/scheme
    color we can't resolve to RGB without instantiating the theme.
    """

    try:
        ft = fill_format.type
    except Exception:
        return None, None
    if ft is None:
        return None, "none"
    # MSO_FILL.SOLID == 1
    if int(ft) == 1:
        try:
            color = fill_format.fore_color
            rgb = color.rgb
            if rgb is not None:
                return rgb_to_hex(rgb), "solid"
        except Exception:
            pass
        return None, "solid"
    if int(ft) == 6:  # MSO_FILL.PICTURE
        return None, "picture"
    if int(ft) == 3:  # MSO_FILL.GRADIENT
        return None, "gradient"
    if int(ft) == 2:  # MSO_FILL.PATTERNED
        return None, "pattern"
    return None, "none"


def slide_background_hex(slide) -> tuple[str, str]:
    """Resolve the effective slide background.

    Returns (hex, source) where source ∈ {slide, layout, master, none}.
    Falls back to white when nothing resolvable.
    """

    for level, getter in [
        ("slide", lambda: slide.background.fill),
        ("layout", lambda: slide.slide_layout.background.fill),
        ("master", lambda: slide.slide_layout.slide_master.background.fill),
    ]:
        try:
            hex_color, ftype = _solid_fill_hex(getter())
            if hex_color:
                return hex_color, level
            if ftype in ("picture", "gradient", "pattern"):
                # Non-solid background — caller will use DEFAULT_BG and mark source.
                return DEFAULT_BG, level
        except Exception:
            continue
    return DEFAULT_BG, "none"


def shape_fill_info(shape) -> dict | None:
    if not hasattr(shape, "fill"):
        return None
    try:
        hex_color, ftype = _solid_fill_hex(shape.fill)
    except Exception:
        return None
    if ftype is None or ftype == "none":
        return None
    info: dict[str, Any] = {"fillType": ftype}
    if hex_color:
        info["fillColor"] = hex_color
    return info


# ── Text run extraction ──────────────────────────────────────────────


def text_runs_from_shape(shape) -> list[dict]:
    """Extract per-run font metrics, falling back to paragraph-level fonts.

    python-pptx returns ``None`` for run-level font.size / font.name / color
    when those properties inherit from the paragraph (or higher). We walk
    one level up to recover the most common case (textbox where size/font
    are set on the paragraph and the run inherits). We do not walk further
    up to layout/master because that would require resolving the placeholder
    type chain — graceful degradation: VLM still has the JPEG.
    """

    runs: list[dict] = []
    if not getattr(shape, "has_text_frame", False):
        return runs
    for para in shape.text_frame.paragraphs:
        # Paragraph-level fallbacks.
        para_size_pt: float | None = None
        try:
            if para.font.size is not None:
                para_size_pt = float(para.font.size.pt)
        except Exception:
            para_size_pt = None
        para_font = None
        try:
            para_font = para.font.name
        except Exception:
            para_font = None
        para_color = None
        try:
            if para.font.color and para.font.color.type:
                para_color = rgb_to_hex(para.font.color.rgb)
        except Exception:
            para_color = None
        para_bold = bool(getattr(para.font, "bold", False))
        para_italic = bool(getattr(para.font, "italic", False))

        for run in para.runs:
            text = run.text or ""
            if not text.strip():
                continue
            font = run.font

            # Run-level with paragraph fallback.
            run_size_pt: float | None = None
            try:
                if font.size is not None:
                    run_size_pt = float(font.size.pt)
            except Exception:
                run_size_pt = None
            font_pt = run_size_pt if run_size_pt is not None else para_size_pt

            font_family = font.name if font.name else para_font

            try:
                run_color = (
                    rgb_to_hex(font.color.rgb)
                    if font.color and font.color.type
                    else None
                )
            except Exception:
                run_color = None
            color = run_color if run_color else para_color

            runs.append(
                {
                    "text": text,
                    "fontPt": font_pt,
                    "fontFamily": font_family,
                    "bold": bool(font.bold) if font.bold is not None else para_bold,
                    "italic": bool(font.italic) if font.italic is not None else para_italic,
                    "color": color,
                }
            )
    return runs


def dominant_text_color(runs: list[dict]) -> str | None:
    """Pick the color of the longest text run with a resolved color."""

    candidates = [r for r in runs if r.get("color")]
    if not candidates:
        return None
    candidates.sort(key=lambda r: len(r.get("text") or ""), reverse=True)
    return candidates[0]["color"]


# ── Zone classification ──────────────────────────────────────────────


def shape_type_label(shape) -> str:
    try:
        st = shape.shape_type
    except Exception:
        return "other"
    if shape.has_text_frame and not getattr(shape, "image", None):
        # placeholder vs plain text — we keep "text" for both; placeholder
        # text presence is captured separately via PLACEHOLDER_TOKENS scan.
        if st == MSO_SHAPE_TYPE.PLACEHOLDER:
            return "placeholder"
        return "text"
    if st == MSO_SHAPE_TYPE.PICTURE:
        return "image"
    if st == MSO_SHAPE_TYPE.GROUP:
        return "group"
    if st == MSO_SHAPE_TYPE.TABLE:
        return "table"
    if st == MSO_SHAPE_TYPE.CHART:
        return "chart"
    if st in (MSO_SHAPE_TYPE.AUTO_SHAPE, MSO_SHAPE_TYPE.FREEFORM):
        return "shape"
    return "other"


def zone_id_for(shape, fallback_idx: int) -> str:
    name = (getattr(shape, "name", "") or "").strip()
    if not name:
        return f"shape-{fallback_idx}"
    # python-pptx auto-names shapes (e.g. "Title 1", "TextBox 5"); collapse
    # whitespace and lowercase to keep ids stable.
    cleaned = re.sub(r"\s+", "-", name.lower())
    cleaned = re.sub(r"[^a-z0-9._-]", "-", cleaned)
    return cleaned or f"shape-{fallback_idx}"


def is_code_zone(zone_id: str, name: str) -> bool:
    haystack = f"{zone_id} {name}".lower()
    return any(hint in haystack for hint in CODE_ZONE_NAME_HINTS)


# ── Image metrics ────────────────────────────────────────────────────


def image_metrics(shape, session_dir: Path) -> dict | None:
    if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
        return None
    try:
        image = shape.image
        blob = image.blob
        ext = image.ext
    except Exception:
        return None

    # We don't have the original disk path — python-pptx exposes the embedded
    # blob. Try to recover the image filename from the relationship if present.
    img_path = None
    try:
        rel_target = shape._element.xpath(".//a:blip/@r:embed")
        # Best-effort: don't fail on missing relationship.
    except Exception:
        rel_target = None

    # Decode dimensions from the blob.
    from io import BytesIO

    try:
        with Image.open(BytesIO(blob)) as im:
            file_w, file_h = im.size
    except Exception:
        return None

    file_aspect = round(file_w / file_h, 4) if file_h else 0.0
    bbox_w_in = emu_to_in(shape.width)
    bbox_h_in = emu_to_in(shape.height)
    zone_aspect = round(bbox_w_in / bbox_h_in, 4) if bbox_h_in else 0.0
    aspect_delta = (
        round(abs(zone_aspect - file_aspect) / file_aspect, 4)
        if file_aspect
        else 0.0
    )

    info: dict[str, Any] = {
        "path": f"<embedded:{ext}>",
        "fileWidthPx": file_w,
        "fileHeightPx": file_h,
        "fileAspectRatio": file_aspect,
        "zoneAspectRatio": zone_aspect,
        "aspectDelta": aspect_delta,
    }
    if rel_target:
        info["path"] = str(rel_target[0]) if hasattr(rel_target[0], "__str__") else info["path"]
    return info


# ── Per-slide assembly ───────────────────────────────────────────────


def compute_margins(
    bbox: dict[str, float], slide_w: float, slide_h: float
) -> dict[str, float]:
    return {
        "left": round(bbox["x"], 4),
        "right": round(slide_w - (bbox["x"] + bbox["w"]), 4),
        "top": round(bbox["y"], 4),
        "bottom": round(slide_h - (bbox["y"] + bbox["h"]), 4),
    }


def detect_overhang(
    bbox: dict[str, float], slide_w: float, slide_h: float
) -> dict[str, float] | None:
    overhang: dict[str, float] = {}
    if bbox["x"] < 0:
        overhang["left"] = round(-bbox["x"], 4)
    if bbox["y"] < 0:
        overhang["top"] = round(-bbox["y"], 4)
    right_over = (bbox["x"] + bbox["w"]) - slide_w
    if right_over > 0:
        overhang["right"] = round(right_over, 4)
    bottom_over = (bbox["y"] + bbox["h"]) - slide_h
    if bottom_over > 0:
        overhang["bottom"] = round(bottom_over, 4)
    return overhang or None


# ── v2 line-wrap estimator (Phase 1, 2026-05-02) ─────────────────────


def _classify_text_kind(font_pt: float | None, role_hint: str | None, font_family: str | None) -> str:
    """Return 'heading' | 'body' | 'mono' for char-width estimation."""
    fam = (font_family or "").lower()
    if any(t in fam for t in ("mono", "consolas", "courier", "menlo", "code")):
        return "mono"
    rh = (role_hint or "").lower()
    if any(tok in rh for tok in HEADING_ROLE_HINT_TOKENS):
        return "heading"
    if font_pt is not None and font_pt >= HEADING_FONT_PT_THRESHOLD:
        return "heading"
    return "body"


def _char_width_em(kind: str) -> float:
    return {
        "heading": CHAR_WIDTH_EM_HEADING,
        "body": CHAR_WIDTH_EM_BODY,
        "mono": CHAR_WIDTH_EM_MONO,
    }.get(kind, CHAR_WIDTH_EM_BODY)


def estimate_wrap(
    text: str,
    font_pt: float | None,
    bbox_w_in: float | None,
    role_hint: str | None,
    font_family: str | None,
) -> tuple[float | None, int | None]:
    """Estimate rendered width (inches) and total line count for a text run.

    Width is based on the longest source line (split on '\\n'); line count
    = max(sourceLineCount, ceil(estimatedWidth / usableWidth)). Returns
    (None, None) when fontPt or bbox unresolvable. See workflow-step8.md
    v2 §5.4.3 (1).
    """
    if font_pt is None or bbox_w_in is None or bbox_w_in <= 0 or not text:
        return None, None
    kind = _classify_text_kind(font_pt, role_hint, font_family)
    char_w_em = _char_width_em(kind)

    source_lines = text.split("\n") or [text]
    longest = max(source_lines, key=len)
    est_width_in = len(longest) * font_pt * char_w_em / 72.0

    # Inset: standard 0.1 in per side, but for narrow zones (e.g. 0.5 in
    # label cells) cap at 10% of zone width to avoid eating most of the
    # usable space and false-flagging "L1"-style tight labels.
    per_side_inset = min(DEFAULT_TEXT_INSET_IN, bbox_w_in * 0.1)
    usable_w_in = max(bbox_w_in - 2 * per_side_inset, 0.01)
    wrap_lines = max(1, int(-(-est_width_in // usable_w_in)))  # ceil
    nonempty_source = len([ln for ln in source_lines if ln.strip()]) or 1
    return round(est_width_in, 3), max(nonempty_source, wrap_lines)


# ── v2 roleHint resolver (Phase 1, 2026-05-02) ───────────────────────


def load_design_role_map(session_dir: Path) -> dict[str, list[str]]:
    """Build slideId → ordered list of zone roles from s05 design JSON.

    The order matches the **build's render order** (s07-build.py iterates
    zones grouped by type: shape → image → formula → text → other), NOT
    the raw design-JSON order. This is what python-pptx exposes when
    iterating slide.shapes after save → only this order aligns with the
    shape index s08 metric uses.

    Used to resolve roleHint by aligning shape order with design JSON zone
    order. Returns empty dict when design JSON missing or unreadable.
    Callers must treat all roleHints as None in that case.
    """
    candidates = [
        "s05-slide-visual-design.json",
        "s06-slide-visual-design.json",
        "design-config.json",
    ]
    design_path: Path | None = None
    for name in candidates:
        cand = session_dir / name
        if cand.is_file():
            design_path = cand
            break
    if design_path is None:
        return {}
    try:
        with open(design_path, encoding="utf-8") as f:
            data = json.load(f)
    except (OSError, json.JSONDecodeError):
        return {}

    # Build's render order, mirrors the four-pass loop in s07-build.py.
    type_order = {"shape": 0, "image": 1, "formula": 2, "text": 3}

    role_map: dict[str, list[str]] = {}
    for sl in data.get("slides", []):
        task_id = sl.get("taskId") or sl.get("slideId")
        if not task_id:
            continue
        zones = (sl.get("layoutSpec") or {}).get("zones") or []
        # Stable sort by render-order group; keep original order within
        # the same type to mirror the build's `for z in zones` inner loops.
        sorted_zones = sorted(
            enumerate(zones),
            key=lambda iz: (type_order.get((iz[1] or {}).get("type"), 99), iz[0]),
        )
        roles = [
            (z.get("role") if isinstance(z, dict) else None) for _, z in sorted_zones
        ]
        role_map[task_id] = roles
    return role_map


def resolve_role_hint(
    slide_id: str, zone_index_zero_based: int, role_map: dict[str, list[str]]
) -> str | None:
    """Look up roleHint by (slideId, zone position). Returns None when out of range."""
    roles = role_map.get(slide_id)
    if not roles:
        return None
    if zone_index_zero_based < 0 or zone_index_zero_based >= len(roles):
        return None
    return roles[zone_index_zero_based] or None


def scan_placeholders(text: str) -> list[tuple[str, str]]:
    """Return list of (matched_token, snippet) tuples found in the text."""
    if not text:
        return []
    hits: list[tuple[str, str]] = []
    lower = text.lower()
    for tok in PLACEHOLDER_TOKENS:
        idx = lower.find(tok.lower())
        if idx == -1:
            continue
        start = max(0, idx - 30)
        end = min(len(text), idx + len(tok) + 30)
        snippet = text[start:end].replace("\n", " ")
        hits.append((tok, snippet))
    return hits


def slide_metrics(
    slide,
    slide_idx: int,
    slide_w: float,
    slide_h: float,
    session_dir: Path,
    role_map: dict[str, list[str]] | None = None,
) -> dict:
    bg_hex, bg_source = slide_background_hex(slide)
    slide_id = f"s{slide_idx:02d}"
    role_map = role_map or {}

    zones: list[dict] = []
    out_of_canvas: list[dict] = []
    placeholder_hits: list[dict] = []

    # Track design-zone position separately from raw shape index — only
    # positionable shapes consume a slot.
    design_zone_idx = -1
    for i, shape in enumerate(slide.shapes, start=1):
        # Skip shapes with no positional info (e.g. some group placeholders).
        if shape.left is None or shape.top is None:
            continue
        design_zone_idx += 1
        role_hint = resolve_role_hint(slide_id, design_zone_idx, role_map)

        zone_id = zone_id_for(shape, i)
        bbox = {
            "x": emu_to_in(shape.left),
            "y": emu_to_in(shape.top),
            "w": emu_to_in(shape.width),
            "h": emu_to_in(shape.height),
        }
        margins = compute_margins(bbox, slide_w, slide_h)
        overhang = detect_overhang(bbox, slide_w, slide_h)
        if overhang:
            out_of_canvas.append({"zoneId": zone_id, "overhangIn": overhang})

        runs = text_runs_from_shape(shape)
        # v2 (Phase 1): augment each run with line-wrap estimate.
        for r in runs:
            est_w, est_lines = estimate_wrap(
                text=r.get("text") or "",
                font_pt=r.get("fontPt"),
                bbox_w_in=bbox.get("w"),
                role_hint=role_hint,
                font_family=r.get("fontFamily"),
            )
            r["estimatedWidthIn"] = est_w
            r["estimatedRenderedLines"] = est_lines
        fill = shape_fill_info(shape)

        # Effective background for contrast: zone fill takes precedence,
        # else slide background.
        eff_bg_hex = (fill or {}).get("fillColor") or bg_hex
        eff_bg_source = "zone-fill" if (fill and fill.get("fillColor")) else f"{bg_source}-bg".replace(
            "slide-bg", "slide-bg"
        )
        if eff_bg_source not in {"zone-fill", "slide-bg", "layout-bg", "master-bg", "unknown"}:
            eff_bg_source = "unknown"

        contrast = None
        fg_hex = dominant_text_color(runs)
        if runs and fg_hex:
            ratio = wcag_contrast(fg_hex, eff_bg_hex)
            contrast = {
                "fgVsBg": ratio,
                "fgRgb": fg_hex,
                "bgRgb": eff_bg_hex,
                "bgSource": eff_bg_source,
            }

        # Placeholder scan over concatenated text in this shape.
        if runs:
            joined = " ".join((r.get("text") or "") for r in runs)
            for tok, snip in scan_placeholders(joined):
                placeholder_hits.append(
                    {"zoneId": zone_id, "match": tok, "snippet": snip}
                )

        zone_entry: dict[str, Any] = {
            "zoneId": zone_id,
            "shapeType": shape_type_label(shape),
            "bboxIn": bbox,
            "marginsIn": margins,
            "fill": fill,
            "textRuns": runs,
            "wcagContrast": contrast,
            "image": image_metrics(shape, session_dir),
            "isCodeZone": is_code_zone(zone_id, getattr(shape, "name", "") or ""),
            "zOrder": i,
            "roleHint": role_hint,
        }
        zones.append(zone_entry)

    return {
        "slideId": slide_id,
        "slideIndex": slide_idx,
        "background": {"fillColor": bg_hex, "source": bg_source},
        "zones": zones,
        "outOfCanvas": out_of_canvas,
        "placeholderHits": placeholder_hits,
    }


# ── Main ─────────────────────────────────────────────────────────────


def main() -> int:
    parser = argparse.ArgumentParser(description=__doc__.split("\n", 1)[0])
    parser.add_argument("session", help="Session directory containing output.pptx")
    parser.add_argument(
        "--output",
        default="output.pptx",
        help="PPTX filename inside the session (default: output.pptx)",
    )
    parser.add_argument(
        "--metrics-out",
        default="s08-zone-metrics.json",
        help="Output metrics filename (default: s08-zone-metrics.json)",
    )
    args = parser.parse_args()

    session = Path(args.session).resolve()
    if not session.is_dir():
        print(f"❌ session dir not found: {session}", file=sys.stderr)
        return 2

    pptx_path = session / args.output
    if not pptx_path.is_file():
        print(f"❌ pptx not found: {pptx_path}", file=sys.stderr)
        return 2

    prs = Presentation(str(pptx_path))
    slide_w_in = emu_to_in(prs.slide_width)
    slide_h_in = emu_to_in(prs.slide_height)

    role_map = load_design_role_map(session)

    slides_out: list[dict] = []
    for idx, slide in enumerate(prs.slides, start=1):
        slides_out.append(slide_metrics(slide, idx, slide_w_in, slide_h_in, session, role_map))

    payload = {
        "slideDimensions": {"widthIn": slide_w_in, "heightIn": slide_h_in},
        "generatedAt": _dt.datetime.now(_dt.timezone.utc).strftime("%Y-%m-%dT%H:%M:%SZ"),
        "sourcePptx": args.output,
        "slides": slides_out,
    }

    out_path = session / args.metrics_out
    with open(out_path, "w", encoding="utf-8") as f:
        json.dump(payload, f, ensure_ascii=False, indent=2)

    n_slides = len(slides_out)
    n_zones = sum(len(s["zones"]) for s in slides_out)
    n_oob = sum(len(s["outOfCanvas"]) for s in slides_out)
    n_ph = sum(len(s["placeholderHits"]) for s in slides_out)
    print(f"✅ wrote {out_path.relative_to(session)}")
    print(f"   slides={n_slides}  zones={n_zones}  outOfCanvas={n_oob}  placeholderHits={n_ph}")
    return 0


if __name__ == "__main__":
    sys.exit(main())
