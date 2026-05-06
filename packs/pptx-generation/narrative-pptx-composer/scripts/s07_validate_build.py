#!/usr/bin/env python3
"""Step 7 checkpoint validator — output.pptx + s05/s06 cross-check.

Cheap structural validator that runs after `s07-build.py` produces
`output.pptx`. Catches three classes of failure that are otherwise only
surfaced by Step 8 visual QA, *after* rendering — costing one extra
render cycle each:

  1. Slide count mismatch — `len(prs.slides) != len(s06.slides)`
     (a builder was missing from the dispatch table, or a duplicate was
     emitted).
  2. Unresolved image path — every plan zone with a non-null `path`
     must point to a file that exists on disk (catches the silent
     image-path resolution bug class flagged in audit item 1.5).
  3. Agenda sentinel still present — no rendered text frame may end
     with the em-dash sentinel `"—"` inside a short label (≤ 12 chars).
     If it does, the 7c agenda page-number patch did not run for that
     slot (audit item 2.5).

Step 8 is still the real downstream gate; this validator only exists to
fail fast before the Step 8 holistic QA subagent burns a full render
budget on a deck with structural defects.

Usage:
  python s07_validate_build.py <session_dir>
  python s07_validate_build.py .                       # from inside session dir
  python s07_validate_build.py . --strict              # exit 1 on any warning
  python s07_validate_build.py . --output other.pptx   # non-default filename
"""

from __future__ import annotations

import argparse
import json
import sys
from pathlib import Path

# Force UTF-8 stdout on Windows so check-mark / cross / warning glyphs
# don't crash the validator on the default cp1252 console.
if sys.platform == "win32":
    try:
        sys.stdout.reconfigure(encoding="utf-8")
        sys.stderr.reconfigure(encoding="utf-8")
    except (AttributeError, OSError):
        pass


# ── Em-dash sentinel ─────────────────────────────────────────────────
# Per SKILL.md § Step 4d: agenda page numbers stay as "—" until 7c
# patches them. The 7c patch replaces them with concrete page references
# (e.g. "pp. 4"). After 7c, no short label should end in the em-dash.
EM_DASH = "\u2014"
SENTINEL_MAX_LEN = 12  # "pp. —" / "—" / "page —" / "—"-only zones


def load_json(path: Path) -> dict:
    with open(path, encoding="utf-8") as f:
        return json.load(f)


# ── Check 1: slide count ─────────────────────────────────────────────

def check_slide_count(prs, content: dict) -> list[str]:
    """Rendered slide count must match s06-slide-content.json entries."""
    rendered = len(prs.slides)
    expected = len(content.get("slides", []))
    if rendered != expected:
        return [
            f"Rendered slide count {rendered} != s06 slide count {expected}. "
            f"Likely cause: a `taskId` is missing from the builder dispatch "
            f"table, or a builder emits more than one slide."
        ]
    return []


# ── Check 2: image paths resolve ─────────────────────────────────────

def _iter_zone_paths(plan: dict):
    """Yield (taskId, zone_role, raw_path) for every zone with a non-null path.

    Walks both `layoutSpec.zones` (body image zones) and
    `layoutSpec.background` (full-bleed background images).
    """
    for slide in plan.get("slides", []):
        tid = slide.get("taskId", "?")
        spec = slide.get("layoutSpec", {}) or {}
        bg = spec.get("background", {}) or {}
        bg_path = bg.get("path")
        if bg_path:
            yield tid, "background", bg_path
        for zone in spec.get("zones", []) or []:
            p = zone.get("path")
            if p:
                yield tid, zone.get("role", "?"), p


def check_image_paths(plan: dict, session_dir: Path) -> list[str]:
    """Every plan zone with a `path` must resolve to an existing file."""
    errors: list[str] = []
    for tid, role, raw in _iter_zone_paths(plan):
        # Normalise Windows backslashes; allow either absolute or session-relative.
        rel = Path(raw.replace("\\", "/"))
        candidate = rel if rel.is_absolute() else (session_dir / rel)
        if not candidate.exists():
            errors.append(f"{tid} / zone '{role}': path '{raw}' does not exist")
    return errors


# ── Check 3: agenda em-dash sentinel cleared ─────────────────────────

def _iter_text_frames(prs):
    """Yield (slide_index_1based, text) for every text frame in the deck.

    Recurses into grouped shapes. Only yields shapes that have text frames;
    pictures, lines, etc. are skipped.
    """
    def _walk(shapes, slide_idx):
        for shape in shapes:
            # GroupShape has .shapes; recurse.
            sub_shapes = getattr(shape, "shapes", None)
            if sub_shapes is not None:
                yield from _walk(sub_shapes, slide_idx)
                continue
            if not getattr(shape, "has_text_frame", False):
                continue
            yield slide_idx, shape.text_frame.text

    for i, slide in enumerate(prs.slides, start=1):
        yield from _walk(slide.shapes, i)


def check_no_dangling_sentinel(prs) -> list[str]:
    """No short text frame may end with the em-dash sentinel after 7c."""
    errors: list[str] = []
    for slide_idx, text in _iter_text_frames(prs):
        if not text:
            continue
        stripped = text.strip()
        if not stripped or EM_DASH not in stripped:
            continue
        # Only flag *short* labels ending in the em-dash — body prose
        # legitimately uses em-dashes mid-sentence everywhere.
        if len(stripped) <= SENTINEL_MAX_LEN and stripped.endswith(EM_DASH):
            errors.append(
                f"Slide {slide_idx}: text frame '{stripped}' still ends with "
                f"em-dash sentinel — Step 7c agenda page-number patch did not "
                f"run, or missed this slot."
            )
    return errors


# ── Check 4: formula artifacts present ────────────────────────────────

def _iter_formula_zones(plan: dict):
    """Yield (taskId, role, zone) for every plan zone with type=='formula'."""
    for slide in plan.get("slides", []):
        tid = slide.get("taskId", "?")
        spec = slide.get("layoutSpec", {}) or {}
        for zone in spec.get("zones", []) or []:
            if zone.get("type") == "formula":
                yield tid, zone.get("role", "?"), zone


def check_formula_artifacts(plan: dict, session_dir: Path) -> list[str]:
    """Every formula zone must have a rendered artifact on disk.

    Phase 1 primary artifact = `<session>/formulas/<taskId>-<role>.png`.
    `.svg` is also accepted (for the Phase 1.5 SVG-injection path).
    """
    errors: list[str] = []
    formulas_dir = session_dir / "formulas"
    for tid, role, _zone in _iter_formula_zones(plan):
        png = formulas_dir / f"{tid}-{role}.png"
        svg = formulas_dir / f"{tid}-{role}.svg"
        if png.exists() and png.stat().st_size > 0:
            continue
        if svg.exists() and svg.stat().st_size > 0:
            continue
        errors.append(
            f"{tid} / formula zone '{role}': no artifact found at "
            f"formulas/{tid}-{role}.png (or .svg). The build script's "
            "render_formulas() pre-pass did not produce this file."
        )
    return errors


# ── Main ─────────────────────────────────────────────────────────────

def main() -> int:
    parser = argparse.ArgumentParser(
        description="Step 7 checkpoint validator for narrative-pptx-composer"
    )
    parser.add_argument(
        "session_dir",
        type=Path,
        help="Path to the session directory containing output.pptx and the JSON artifacts",
    )
    parser.add_argument(
        "--output",
        default="output.pptx",
        help="Output .pptx filename inside session_dir (default: output.pptx)",
    )
    parser.add_argument(
        "--strict",
        action="store_true",
        help="Treat warnings as errors (exit 1 on any issue)",
    )
    args = parser.parse_args()

    session = args.session_dir.resolve()
    output_path = session / args.output

    if not output_path.exists():
        print(f"FAIL: {args.output} not found in {session}")
        return 1

    # Load required JSON artifacts
    required = {
        "s05-slide-visual-design.json": "plan",
        "s06-slide-content.json": "content",
    }
    artifacts: dict[str, dict] = {}
    for filename, key in required.items():
        path = session / filename
        if not path.exists():
            print(f"FAIL: {filename} not found in {session}")
            return 1
        artifacts[key] = load_json(path)

    plan = artifacts["plan"]
    content = artifacts["content"]

    try:
        from pptx import Presentation
    except ImportError:
        print("FAIL: python-pptx is not installed; run `pip install python-pptx`")
        return 1

    try:
        prs = Presentation(str(output_path))
    except Exception as exc:  # pptx.exc.PackageNotFoundError, BadZipFile, …
        print(f"FAIL: could not open {args.output} as a .pptx package: {exc}")
        return 1

    print(f"Validating Step 7 checkpoint in: {session.name}")
    print(f"  Output: {args.output}")
    print(f"  Plan tasks: {len(plan.get('slides', []))}")
    print(f"  Content entries: {len(content.get('slides', []))}")
    print(f"  Rendered slides: {len(prs.slides)}")
    print()

    all_errors: list[str] = []
    all_warnings: list[str] = []

    # 1. Slide count
    errors = check_slide_count(prs, content)
    if errors:
        print("✗ Slide count")
        for e in errors:
            print(f"  ERROR: {e}")
        all_errors.extend(errors)
    else:
        print(f"✓ Slide count — {len(prs.slides)}/{len(content['slides'])} match")

    # 2. Image paths
    errors = check_image_paths(plan, session)
    if errors:
        print("✗ Image paths")
        for e in errors:
            print(f"  ERROR: {e}")
        all_errors.extend(errors)
    else:
        path_count = sum(1 for _ in _iter_zone_paths(plan))
        if path_count:
            print(f"✓ Image paths — {path_count} plan zone(s) all resolve")
        else:
            print("ℹ Image paths — no image zones in plan")

    # 3. Agenda em-dash sentinel
    errors = check_no_dangling_sentinel(prs)
    if errors:
        print("✗ Agenda sentinel")
        for e in errors:
            print(f"  ERROR: {e}")
        all_errors.extend(errors)
    else:
        print("✓ Agenda sentinel — no short labels end with em-dash (7c patched)")

    # 4. Formula artifacts
    errors = check_formula_artifacts(plan, session)
    formula_count = sum(1 for _ in _iter_formula_zones(plan))
    if errors:
        print("✗ Formula artifacts")
        for e in errors:
            print(f"  ERROR: {e}")
        all_errors.extend(errors)
    elif formula_count:
        print(f"✓ Formula artifacts — {formula_count} formula(s) rendered")

    # Summary
    total_issues = len(all_errors) + (len(all_warnings) if args.strict else 0)
    if total_issues == 0:
        print("═══ PASS — Step 7 checkpoint validated ═══")
        return 0
    print(f"═══ FAIL — {len(all_errors)} errors, {len(all_warnings)} warnings ═══")
    return 1


if __name__ == "__main__":
    sys.exit(main())
