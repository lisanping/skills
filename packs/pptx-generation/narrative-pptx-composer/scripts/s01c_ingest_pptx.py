"""Extract structured content and embedded images from a .pptx file.

Produces a content-digest JSON (s01c-content-digest.json format) that
the narrative-pptx-composer workflow can consume directly.  Extracts:

  - Per-slide text (title, body, table cells, grouped shapes)
  - Embedded images → copied to session images/ directory
  - Speaker notes
  - Slide-level metadata (layout name, slide number)
  - Data points and named entities (pre-glossary seeds)

Usage:
    python s01c_ingest_pptx.py <input.pptx> <session_dir>

Outputs (inside <session_dir>):
    s01c-content-digest.json   — structured content digest
    images/                    — extracted embedded images

Requires: python-pptx, Pillow
"""

from __future__ import annotations

import argparse
import hashlib
import json
import re
import sys
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


# ── Helpers ────────────────────────────────────────────────────────


def _slug(text: str, max_len: int = 40) -> str:
    """Lowercase ASCII slug for filenames."""
    s = re.sub(r"[^a-z0-9]+", "-", text.lower()).strip("-")
    return s[:max_len] if s else "untitled"


def _short_hash(data: bytes, length: int = 8) -> str:
    return hashlib.sha256(data).hexdigest()[:length]


def _shape_text(shape) -> str:
    """Return the full text of a shape's text_frame, or ''."""
    if not shape.has_text_frame:
        return ""
    return "\n".join(p.text for p in shape.text_frame.paragraphs if p.text.strip())


def _table_to_text(table) -> str:
    """Render a pptx Table as a markdown-style table string."""
    rows = []
    for row in table.rows:
        cells = [cell.text.strip() for cell in row.cells]
        rows.append("| " + " | ".join(cells) + " |")
    if len(rows) >= 2:
        # Insert header separator after first row
        col_count = len(table.columns)
        rows.insert(1, "| " + " | ".join(["---"] * col_count) + " |")
    return "\n".join(rows)


def _extract_data_points(text: str) -> list[str]:
    """Heuristic: pull numbers with units / percentages / currencies."""
    patterns = [
        r"[\$€¥£]\s?\d[\d,]*\.?\d*[BMKbmk]?",       # currencies
        r"\d[\d,]*\.?\d*\s?%",                         # percentages
        r"\d[\d,]*\.?\d*\s?(?:billion|million|thousand|[BMK])\b",  # big numbers
    ]
    combined = "|".join(f"({p})" for p in patterns)
    return [m.group(0).strip() for m in re.finditer(combined, text, re.IGNORECASE)]


def _guess_role(title: str, body: str) -> str:
    """Heuristic suggested-role for a slide section."""
    combined = (title + " " + body).lower()
    data_signals = ["chart", "graph", "table", "metric", "data", "kpi",
                    "revenue", "growth", "budget", "forecast", "pipeline",
                    "%", "$", "€", "¥"]
    data_hits = sum(1 for s in data_signals if s in combined)
    if data_hits >= 2:
        return "content_data"
    if data_hits == 1 and len(body) < 200:
        return "content_data"
    if len(body) > 400:
        return "content_text"
    return "content_mixed"


# ── Image extraction ───────────────────────────────────────────────


def _extract_images(prs: Presentation, slide, slide_idx: int,
                    images_dir: Path,
                    seen_hashes: dict[str, str] | None = None,
                    ) -> list[dict[str, Any]]:
    """Extract embedded images from a slide. Returns mediaAsset entries.

    Args:
        seen_hashes: shared dict mapping blob-hash → relative path.
            Prevents writing the same image bytes more than once.
    """
    if seen_hashes is None:
        seen_hashes = {}
    assets: list[dict[str, Any]] = []
    img_counter = 0

    for shape in slide.shapes:
        blob = None
        content_type = None

        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            try:
                image = shape.image
                blob = image.blob
                content_type = image.content_type
            except ValueError:
                # Linked image (not embedded) — skip
                pass
        elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            # Walk group children for pictures
            for child in shape.shapes:
                if child.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    try:
                        image = child.image
                        blob = image.blob
                        content_type = image.content_type
                    except ValueError:
                        pass
                    break

        if blob is None:
            continue

        img_counter += 1
        h = _short_hash(blob)

        # Deduplicate: reuse path if same content was already written
        if h in seen_hashes:
            rel_path = seen_hashes[h]
        else:
            ext = {
                "image/png": ".png",
                "image/jpeg": ".jpg",
                "image/gif": ".gif",
                "image/bmp": ".bmp",
                "image/tiff": ".tiff",
                "image/svg+xml": ".svg",
                "image/x-wmf": ".wmf",
                "image/x-emf": ".emf",
            }.get(content_type, ".png")

            fname = f"slide{slide_idx + 1:02d}_img{img_counter:02d}_{h}{ext}"
            out_path = images_dir / fname
            out_path.write_bytes(blob)
            rel_path = f"images/{fname}"
            seen_hashes[h] = rel_path

        alt_text = ""
        if hasattr(shape, "name"):
            alt_text = shape.name

        assets.append({
            "type": "image",
            "ref": rel_path,
            "alt": alt_text,
            "context": f"Embedded image from slide {slide_idx + 1}",
            "sourceAnchor": f"pptx-slide-{slide_idx + 1}",
        })

    return assets


# ── Per-slide extraction ───────────────────────────────────────────


def _extract_slide(prs: Presentation, slide, slide_idx: int,
                   images_dir: Path,
                   seen_hashes: dict[str, str] | None = None,
                   ) -> dict[str, Any]:
    """Extract structured content from a single slide."""
    layout_name = ""
    if slide.slide_layout:
        layout_name = slide.slide_layout.name or ""

    title_text = ""
    body_parts: list[str] = []
    table_parts: list[str] = []

    # First pass: find title via placeholder idx 0 or shapes.title
    title_shape = None
    if slide.shapes.title is not None:
        title_shape = slide.shapes.title
        title_text = title_shape.text.strip() if title_shape.text else ""

    for shape in slide.shapes:
        if shape is title_shape:
            continue

        if not title_text and hasattr(shape, "is_placeholder") and shape.is_placeholder:
            phf = shape.placeholder_format
            if phf.idx == 0:
                title_text = _shape_text(shape)
                continue

        if shape.has_table:
            table_parts.append(_table_to_text(shape.table))
            continue

        txt = _shape_text(shape)
        if txt:
            body_parts.append(txt)

    # Speaker notes
    notes = ""
    if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
        notes = slide.notes_slide.notes_text_frame.text.strip()

    # Combine body
    body = "\n\n".join(body_parts)
    if table_parts:
        body += ("\n\n" if body else "") + "\n\n".join(table_parts)

    full_text = f"{title_text}\n{body}" if title_text else body
    data_points = _extract_data_points(full_text)
    role = _guess_role(title_text, body)

    # Images
    media_assets = _extract_images(prs, slide, slide_idx, images_dir, seen_hashes)

    anchor = f"pptx-slide-{slide_idx + 1}"

    section: dict[str, Any] = {
        "heading": title_text or f"Slide {slide_idx + 1}",
        "level": 1,
        "content": body,
        "dataPoints": data_points,
        "suggestedRole": role,
        "anchor": anchor,
        "layoutName": layout_name,
        "mediaAssets": media_assets,
    }
    if notes:
        section["speakerNotes"] = notes

    return section


# ── Main ───────────────────────────────────────────────────────────


def ingest(pptx_path: Path, session_dir: Path) -> dict[str, Any]:
    """Parse a .pptx and produce a content-digest dict."""
    prs = Presentation(str(pptx_path))
    images_dir = session_dir / "images"
    images_dir.mkdir(parents=True, exist_ok=True)

    sections: list[dict[str, Any]] = []
    all_media: list[dict[str, Any]] = []
    all_data_points: list[str] = []
    total_chars = 0
    seen_hashes: dict[str, str] = {}

    for idx, slide in enumerate(prs.slides):
        sec = _extract_slide(prs, slide, idx, images_dir, seen_hashes)
        sections.append(sec)
        total_chars += len(sec.get("content", ""))
        all_data_points.extend(sec.get("dataPoints", []))
        for ma in sec.get("mediaAssets", []):
            all_media.append({**ma, "sourceFile": pptx_path.name})

    # Build pre-glossary from data points
    pre_glossary: dict[str, str] = {}
    for i, dp in enumerate(dict.fromkeys(all_data_points)):
        key = _slug(dp, 30).replace("-", "_") or f"dp_{i}"
        pre_glossary[key] = dp

    # Dimensions
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    dims = {
        "widthEmu": slide_width,
        "heightEmu": slide_height,
        "widthInches": round(slide_width / 914400, 2),
        "heightInches": round(slide_height / 914400, 2),
    }

    digest: dict[str, Any] = {
        "sources": [
            {
                "file": pptx_path.name,
                "format": "pptx",
                "slideCount": len(prs.slides),
                "dimensions": dims,
                "sections": sections,
                "totalChars": total_chars,
            }
        ],
        "mergedSections": [
            {**sec, "sourceFile": pptx_path.name}
            for sec in sections
        ],
        "mediaAssets": all_media,
        "preGlossary": pre_glossary,
    }
    return digest


def main() -> None:
    parser = argparse.ArgumentParser(
        description="Extract structured content from a .pptx into content-digest JSON."
    )
    parser.add_argument("input", help="Input .pptx file")
    parser.add_argument("session_dir", help="Session directory for outputs")
    args = parser.parse_args()

    pptx_path = Path(args.input).resolve()
    session_dir = Path(args.session_dir).resolve()

    if not pptx_path.exists():
        print(f"Error: {pptx_path} not found", file=sys.stderr)
        sys.exit(1)
    if not pptx_path.suffix.lower() == ".pptx":
        if pptx_path.suffix.lower() == ".potx":
            print("Error: .potx is a template format — use the pptx-profiler "
                  "skill instead. This script ingests .pptx content files.",
                  file=sys.stderr)
        else:
            print(f"Error: expected .pptx, got {pptx_path.suffix}", file=sys.stderr)
        sys.exit(1)

    session_dir.mkdir(parents=True, exist_ok=True)

    digest = ingest(pptx_path, session_dir)

    out_path = session_dir / "s01c-content-digest.json"
    out_path.write_text(json.dumps(digest, indent=2, ensure_ascii=False), encoding="utf-8")

    slide_count = digest["sources"][0]["slideCount"]
    image_count = len(digest["mediaAssets"])
    section_count = len(digest["mergedSections"])
    print(f"Ingested {pptx_path.name}: {slide_count} slides, "
          f"{section_count} sections, {image_count} images extracted")
    print(f"Output: {out_path}")


if __name__ == "__main__":
    main()
