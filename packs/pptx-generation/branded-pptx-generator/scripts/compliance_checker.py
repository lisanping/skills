"""
Compliance Checker — Post-hoc brand compliance validation.

Reads a template_profile.json and a generated .pptx, then checks
Locked-tier rules (colors, fonts, slide size) against the template's theme.

Usage:
    python compliance_checker.py output.pptx --profile template_profile.json
    python compliance_checker.py output.pptx --profile template_profile.json --strict

Requires: python-pptx (pip install python-pptx)
"""

import argparse
import json
import re
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.dml import MSO_THEME_COLOR
except ImportError:
    print("ERROR: python-pptx required. Install: pip install python-pptx", file=sys.stderr)
    sys.exit(1)

from lxml import etree

# Namespace map for XML queries
_NSMAP = {
    'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
    'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
}


# ── Color utilities ──────────────────────────────────────

def normalize_hex(color: str) -> str:
    """Normalize a hex color to uppercase without '#'."""
    return color.lstrip("#").upper()


def rgb_to_hex(rgb: RGBColor) -> str:
    return f"{rgb.red:02X}{rgb.green:02X}{rgb.blue:02X}"


def relative_luminance(hex_color: str) -> float:
    """WCAG relative luminance from hex color."""
    r, g, b = int(hex_color[:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16)
    vals = []
    for c in (r, g, b):
        s = c / 255.0
        vals.append(s / 12.92 if s <= 0.03928 else ((s + 0.055) / 1.055) ** 2.4)
    return 0.2126 * vals[0] + 0.7152 * vals[1] + 0.0722 * vals[2]


def contrast_ratio(hex1: str, hex2: str) -> float:
    """WCAG contrast ratio between two hex colors."""
    l1 = relative_luminance(normalize_hex(hex1))
    l2 = relative_luminance(normalize_hex(hex2))
    lighter, darker = max(l1, l2), min(l1, l2)
    return (lighter + 0.05) / (darker + 0.05)


# ── Profile loading ──────────────────────────────────────

def load_profile(path: str) -> dict:
    with open(path, "r", encoding="utf-8") as f:
        return json.load(f)


def get_allowed_colors(profile: dict) -> set[str]:
    """Extract all allowed colors from the theme scheme."""
    scheme = profile.get("identity", {}).get("colors", {}).get("scheme", {})
    colors = set()
    for v in scheme.values():
        if v:
            colors.add(normalize_hex(v))
    return colors


def get_allowed_fonts(profile: dict) -> set[str]:
    """Extract all allowed font families from the profile."""
    fonts_section = profile.get("identity", {}).get("fonts", {})
    fonts = set()
    for group in ("major", "minor"):
        g = fonts_section.get(group, {})
        for script in ("latin", "ea", "cs"):
            f = g.get(script)
            if f:
                fonts.add(f)
    return fonts


# Map MSO_THEME_COLOR enum → profile scheme key
_THEME_TO_SCHEME = {
    MSO_THEME_COLOR.DARK_1: "dk1",
    MSO_THEME_COLOR.LIGHT_1: "lt1",
    MSO_THEME_COLOR.DARK_2: "dk2",
    MSO_THEME_COLOR.LIGHT_2: "lt2",
    MSO_THEME_COLOR.ACCENT_1: "accent1",
    MSO_THEME_COLOR.ACCENT_2: "accent2",
    MSO_THEME_COLOR.ACCENT_3: "accent3",
    MSO_THEME_COLOR.ACCENT_4: "accent4",
    MSO_THEME_COLOR.ACCENT_5: "accent5",
    MSO_THEME_COLOR.ACCENT_6: "accent6",
    MSO_THEME_COLOR.HYPERLINK: "hlink",
    MSO_THEME_COLOR.FOLLOWED_HYPERLINK: "folHlink",
    MSO_THEME_COLOR.BACKGROUND_1: "lt1",
    MSO_THEME_COLOR.BACKGROUND_2: "lt2",
    MSO_THEME_COLOR.TEXT_1: "dk1",
    MSO_THEME_COLOR.TEXT_2: "dk2",
}

# Map XML schemeClr val → profile scheme key
_XML_SCHEME_MAP = {
    "dk1": "dk1", "dk2": "dk2", "lt1": "lt1", "lt2": "lt2",
    "accent1": "accent1", "accent2": "accent2", "accent3": "accent3",
    "accent4": "accent4", "accent5": "accent5", "accent6": "accent6",
    "hlink": "hlink", "folHlink": "folHlink",
    "bg1": "lt1", "bg2": "lt2", "tx1": "dk1", "tx2": "dk2",
}


def build_scheme_rgb_map(profile: dict) -> dict[str, str]:
    """Build {scheme_key: 'RRGGBB'} from the profile's color scheme."""
    scheme = profile.get("identity", {}).get("colors", {}).get("scheme", {})
    return {k: normalize_hex(v) for k, v in scheme.items() if v}


def resolve_theme_color(theme_color_enum, scheme_map: dict[str, str]) -> str | None:
    """Resolve an MSO_THEME_COLOR enum to a hex RGB string via profile scheme."""
    key = _THEME_TO_SCHEME.get(theme_color_enum)
    if key:
        return scheme_map.get(key)
    return None


# ── PPTX scanning ────────────────────────────────────────

def scan_pptx(pptx_path: str, scheme_map: dict[str, str]) -> dict:
    """Scan a .pptx and collect all used colors, fonts, and text elements."""
    prs = Presentation(pptx_path)
    colors_used: set[str] = set()
    fonts_used: set[str] = set()
    text_elements: list[dict] = []  # for WCAG checks

    for slide_idx, slide in enumerate(prs.slides):
        bg_color = _get_slide_bg(slide, scheme_map)

        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        # Font
                        if run.font.name:
                            fonts_used.add(run.font.name)
                        # Color — resolve both RGB and theme colors
                        fg_hex = None
                        is_theme = False
                        try:
                            if run.font.color and run.font.color.rgb:
                                fg_hex = rgb_to_hex(run.font.color.rgb)
                                colors_used.add(fg_hex)
                        except AttributeError:
                            # Theme/scheme colors: compliant for locked-tier
                            # but still need WCAG contrast check
                            try:
                                tc = run.font.color.theme_color
                                if tc is not None:
                                    fg_hex = resolve_theme_color(tc, scheme_map)
                                    is_theme = True
                            except (AttributeError, TypeError):
                                pass

                        if fg_hex:
                            text_elements.append({
                                "slide": slide_idx + 1,
                                "text": run.text[:50],
                                "fg": fg_hex,
                                "bg": bg_color,
                                "size_pt": run.font.size.pt if run.font.size else None,
                                "is_theme": is_theme,
                            })

            # Shape fill colors
            if hasattr(shape, "fill"):
                fill = shape.fill
                if fill.type is not None:
                    try:
                        if fill.fore_color and fill.fore_color.rgb:
                            colors_used.add(rgb_to_hex(fill.fore_color.rgb))
                    except (AttributeError, TypeError):
                        pass

    return {
        "colors_used": colors_used,
        "fonts_used": fonts_used,
        "text_elements": text_elements,
    }


def _get_slide_bg(slide, scheme_map: dict[str, str]) -> str:
    """Best-effort slide background color extraction.

    Resolution order: slide bg → layout bg → master bg → profile lt1 → white.
    Handles both RGB fills and theme-color (schemeClr) fills.
    """
    # Try the slide's own background
    resolved = _resolve_bg_from_element(slide._element, scheme_map)
    if resolved:
        return resolved

    # Fall back to slide layout
    try:
        layout = slide.slide_layout
        resolved = _resolve_bg_from_element(layout._element, scheme_map)
        if resolved:
            return resolved
    except (AttributeError, TypeError):
        pass

    # Fall back to slide master
    try:
        master = slide.slide_layout.slide_master
        resolved = _resolve_bg_from_element(master._element, scheme_map)
        if resolved:
            return resolved
    except (AttributeError, TypeError):
        pass

    # Final fallback: profile lt1 (typical light background) or white
    return scheme_map.get("lt1", "FFFFFF")


def _resolve_bg_from_element(element, scheme_map: dict[str, str]) -> str | None:
    """Extract background color from a slide/layout/master XML element."""
    # Look for <p:bg> → <p:bgPr> → <a:solidFill>
    bg = element.findall('.//p:bg/p:bgPr/a:solidFill', _NSMAP)
    if not bg:
        # Also check <p:bg> → <p:bgRef> (master-style reference)
        bg = element.findall('.//p:bg/p:bgRef', _NSMAP)
    for fill_el in bg:
        # Check for <a:srgbClr val="RRGGBB"/>
        srgb = fill_el.find('.//a:srgbClr', _NSMAP)
        if srgb is not None:
            val = srgb.get('val')
            if val:
                return normalize_hex(val)
        # Check for <a:schemeClr val="accent2"/>
        scheme_clr = fill_el.find('.//a:schemeClr', _NSMAP)
        if scheme_clr is not None:
            val = scheme_clr.get('val')
            key = _XML_SCHEME_MAP.get(val)
            if key:
                return scheme_map.get(key)
    return None


# ── Compliance checks ────────────────────────────────────

def check_locked_colors(colors_used: set[str], allowed: set[str]) -> list[dict]:
    violations = []
    for color in sorted(colors_used):
        if color not in allowed:
            violations.append({
                "tier": "locked",
                "element": "color",
                "severity": "critical",
                "expected": f"One of theme colors",
                "actual": f"#{color}",
            })
    return violations


def check_locked_fonts(fonts_used: set[str], allowed: set[str]) -> list[dict]:
    violations = []
    for font in sorted(fonts_used):
        if font not in allowed:
            violations.append({
                "tier": "locked",
                "element": "font",
                "severity": "critical",
                "expected": f"One of: {', '.join(sorted(allowed))}",
                "actual": font,
            })
    return violations


def check_locked_slide_size(pptx_path: str, profile: dict) -> list[dict]:
    """Check that the generated PPTX slide size matches the template profile."""
    expected = profile.get("identity", {}).get("slide_size", {})
    if not expected:
        return []
    prs = Presentation(pptx_path)
    violations = []
    expected_w = expected.get("width_pt")
    expected_h = expected.get("height_pt")
    if expected_w is not None:
        actual_w = prs.slide_width.pt
        if abs(actual_w - expected_w) > 0.5:
            violations.append({
                "tier": "locked",
                "element": "slide_size_width",
                "severity": "critical",
                "expected": f"{expected_w}pt",
                "actual": f"{actual_w}pt",
            })
    if expected_h is not None:
        actual_h = prs.slide_height.pt
        if abs(actual_h - expected_h) > 0.5:
            violations.append({
                "tier": "locked",
                "element": "slide_size_height",
                "severity": "critical",
                "expected": f"{expected_h}pt",
                "actual": f"{actual_h}pt",
            })
    return violations


def check_wcag(text_elements: list[dict]) -> list[dict]:
    issues = []
    for elem in text_elements:
        fg = elem.get("fg")
        bg = elem.get("bg", "FFFFFF")
        if not fg or not bg:
            continue
        ratio = contrast_ratio(fg, bg)
        size = elem.get("size_pt") or 14
        required = 3.0 if size >= 24 else 4.5
        if ratio < required:
            issues.append({
                "slide": elem["slide"],
                "text": elem["text"],
                "contrast_ratio": round(ratio, 2),
                "required_ratio": required,
                "fg": f"#{fg}",
                "bg": f"#{bg}",
            })
    return issues


# ── Main ─────────────────────────────────────────────────

def main():
    parser = argparse.ArgumentParser(
        description="Check brand compliance of a generated .pptx against template profile",
    )
    parser.add_argument("pptx", help="Path to the generated .pptx file")
    parser.add_argument(
        "--profile", "-p",
        required=True,
        help="Path to template_profile.json",
    )
    parser.add_argument(
        "--output", "-o",
        help="Output JSON report path (default: print to stdout)",
    )
    parser.add_argument(
        "--strict",
        action="store_true",
        help="Exit with code 1 if any locked violations found",
    )

    args = parser.parse_args()

    for path, label in [(args.pptx, "PPTX"), (args.profile, "Profile")]:
        if not Path(path).exists():
            print(f"ERROR: {label} not found: {path}", file=sys.stderr)
            sys.exit(1)

    profile = load_profile(args.profile)
    allowed_colors = get_allowed_colors(profile)
    allowed_fonts = get_allowed_fonts(profile)
    scheme_map = build_scheme_rgb_map(profile)

    print(f"Profile: {len(allowed_colors)} theme colors, {len(allowed_fonts)} font families")
    print(f"Scanning: {args.pptx}")

    scan = scan_pptx(args.pptx, scheme_map)
    print(f"Found: {len(scan['colors_used'])} colors, {len(scan['fonts_used'])} fonts, {len(scan['text_elements'])} text runs")

    # Run checks
    color_violations = check_locked_colors(scan["colors_used"], allowed_colors)
    font_violations = check_locked_fonts(scan["fonts_used"], allowed_fonts)
    slide_size_violations = check_locked_slide_size(args.pptx, profile)
    wcag_issues = check_wcag(scan["text_elements"])

    all_locked = color_violations + font_violations + slide_size_violations
    report = {
        "pptx_file": str(Path(args.pptx).name),
        "profile_file": str(Path(args.profile).name),
        "locked_violations": all_locked,
        "wcag_issues": wcag_issues,
        "verdict": "pass" if len(all_locked) == 0 else "fail",
        "summary": {
            "colors_checked": len(scan["colors_used"]),
            "fonts_checked": len(scan["fonts_used"]),
            "color_violations": len(color_violations),
            "font_violations": len(font_violations),
            "slide_size_violations": len(slide_size_violations),
            "wcag_issues": len(wcag_issues),
        },
    }

    # Print summary
    print(f"\n{'='*50}")
    print(f"VERDICT: {report['verdict'].upper()}")
    print(f"{'='*50}")

    if color_violations:
        print(f"\n🔒 Color violations ({len(color_violations)}):")
        for v in color_violations:
            print(f"  ✗ {v['actual']} — not in theme scheme")

    if font_violations:
        print(f"\n🔒 Font violations ({len(font_violations)}):")
        for v in font_violations:
            print(f"  ✗ {v['actual']} — expected: {v['expected']}")

    if slide_size_violations:
        print(f"\n🔒 Slide size violations ({len(slide_size_violations)}):")
        for v in slide_size_violations:
            print(f"  ✗ {v['element']}: {v['actual']} — expected: {v['expected']}")

    if wcag_issues:
        print(f"\n⚠ WCAG contrast issues ({len(wcag_issues)}):")
        for w in wcag_issues:
            print(
                f"  Slide {w['slide']}: \"{w['text']}\" — "
                f"ratio {w['contrast_ratio']} < {w['required_ratio']} "
                f"(fg={w['fg']}, bg={w['bg']})"
            )

    if not all_locked and not wcag_issues:
        print("\n✓ All checks passed")

    # Output JSON if requested
    if args.output:
        out_path = Path(args.output)
        out_path.parent.mkdir(parents=True, exist_ok=True)
        out_path.write_text(json.dumps(report, indent=2, ensure_ascii=False))
        print(f"\n✓ Report written to {out_path}")

    if args.strict and report["verdict"] == "fail":
        sys.exit(1)


if __name__ == "__main__":
    main()
